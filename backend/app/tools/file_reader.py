import os
import csv
import json
from typing import Optional

from app.tools.excel_analyzer import ExcelAnalyzer


class FileReader:
    """
    Reads and extracts text content from various file formats.
    Supports: PDF, DOCX, XLSX/XLS, PPTX, CSV, JSON, XML, TXT, MD, images, and code files.
    """

    async def extract_text(
        self,
        file_path: str,
        file_type: str,
        original_name: str = "",
        max_length: int = 50000,
    ) -> str:
        """
        Extract text content from a file.

        Args:
            file_path: Path to the file
            file_type: Type identifier (pdf, docx, xlsx, etc.)
            original_name: Original filename for display
            max_length: Maximum characters to return

        Returns:
            Extracted text content
        """
        if not os.path.exists(file_path):
            return f"[文件不存在: {file_path}]"

        try:
            text = await self._extract_by_type(file_path, file_type, original_name)
            if text and len(text) > max_length:
                text = text[:max_length] + f"\n\n[... 内容过长，已截断，显示前 {max_length} 字符 ...]"
            return text or f"[无法从 {original_name or file_type} 文件提取内容]"
        except Exception as e:
            return f"[文件内容提取失败 ({file_type}): {str(e)}]"

    async def _extract_by_type(self, file_path: str, file_type: str, original_name: str) -> str:
        """Route extraction to appropriate method based on file type."""
        file_type = file_type.lower()

        if file_type == "pdf":
            return await self._extract_pdf(file_path)
        elif file_type in ("docx", "doc"):
            return await self._extract_docx(file_path)
        elif file_type in ("xlsx", "xls"):
            return await self._extract_xlsx(file_path)
        elif file_type == "csv":
            return await self._extract_csv(file_path)
        elif file_type == "pptx":
            return await self._extract_pptx(file_path)
        elif file_type == "json":
            return await self._extract_json(file_path)
        elif file_type in ("xml", "html", ".html"):
            return await self._extract_xml(file_path)
        elif file_type in ("txt", "markdown", "md"):
            return await self._extract_text_file(file_path)
        elif file_type == "image":
            return f"[图片文件: {original_name or os.path.basename(file_path)}] （图片内容需通过视觉分析处理）"
        elif file_type in (
            "python", "javascript", "typescript", "java", "cpp", "c",
            "go", "rust", "sql", "shell", "yaml", "toml", "ini", "config"
        ):
            return await self._extract_text_file(file_path)
        else:
            # Try as text file
            try:
                return await self._extract_text_file(file_path)
            except Exception:
                return f"[二进制文件: {original_name or os.path.basename(file_path)}] （无法提取文本内容）"

    async def _extract_pdf(self, file_path: str) -> str:
        """Extract text + tables from PDF.

        Strategy:
          1. pdfplumber (首选) — 精确表格提取 + 文字层
          2. PyPDF2 (fallback) — 仅文字层
        """
        # ── pdfplumber path ─────────────────────────────────────────────
        try:
            import pdfplumber
            text_parts = []
            with pdfplumber.open(file_path) as pdf:
                total = len(pdf.pages)
                for i, page in enumerate(pdf.pages):
                    page_header = f"[第 {i+1}/{total} 页]"
                    page_texts: list[str] = []

                    # 正文
                    raw = page.extract_text(x_tolerance=3, y_tolerance=3)
                    if raw and raw.strip():
                        page_texts.append(raw.strip())

                    # 表格（结构化保留）
                    tables = page.extract_tables()
                    for t_idx, table in enumerate(tables):
                        if not table:
                            continue
                        rows_md = []
                        for r_idx, row in enumerate(table):
                            cells = [str(c).strip() if c else "" for c in row]
                            rows_md.append("| " + " | ".join(cells) + " |")
                            if r_idx == 0:
                                rows_md.append("| " + " | ".join(["---"] * len(cells)) + " |")
                        page_texts.append(f"\n[表格 {t_idx+1}]\n" + "\n".join(rows_md))

                    if page_texts:
                        text_parts.append(page_header + "\n" + "\n".join(page_texts))
                    else:
                        text_parts.append(page_header + "\n[本页无可提取文字，可能为扫描件]")

            return "\n\n".join(text_parts) if text_parts else "PDF文件无可提取文本"

        except ImportError:
            pass  # fall through to PyPDF2
        except Exception:
            pass

        # ── PyPDF2 fallback ─────────────────────────────────────────────
        try:
            import PyPDF2
            text_parts = []
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                total_pages = len(reader.pages)
                for i, page in enumerate(reader.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text:
                            text_parts.append(f"[第 {i+1}/{total_pages} 页]\n{page_text}")
                    except Exception:
                        text_parts.append(f"[第 {i+1}/{total_pages} 页 - 提取失败]")
            return "\n\n".join(text_parts) if text_parts else "PDF文件无可提取文本"
        except ImportError:
            return "[pdfplumber/PyPDF2 均未安装，无法提取PDF内容]"
        except Exception as e:
            return f"[PDF提取失败: {str(e)}]"

    async def _extract_docx(self, file_path: str) -> str:
        """Extract text from DOCX using python-docx."""
        try:
            from docx import Document
            doc = Document(file_path)
            parts = []

            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    # Add heading markers
                    if para.style.name.startswith("Heading"):
                        level = para.style.name.replace("Heading ", "")
                        parts.append(f"\n{'#' * int(level) if level.isdigit() else '#'} {para.text}")
                    else:
                        parts.append(para.text)

            # Extract tables
            for i, table in enumerate(doc.tables):
                parts.append(f"\n[表格 {i+1}]")
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        parts.append(row_text)

            return "\n".join(parts)
        except ImportError:
            return "[python-docx未安装，无法提取DOCX内容]"
        except Exception as e:
            return f"[DOCX提取失败: {str(e)}]"

    async def _extract_xlsx(self, file_path: str) -> str:
        """Extract text from XLSX/XLS using openpyxl or xlrd."""
        analyzer = ExcelAnalyzer()
        # Try openpyxl first (for .xlsx)
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            parts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                parts.append(f"\n[工作表: {sheet_name}]")
                row_count = 0
                for row in ws.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                        parts.append(row_text)
                        row_count += 1
                        if row_count > 500:  # Limit rows per sheet
                            parts.append(f"[... 仅显示前500行，共更多行 ...]")
                            break
            wb.close()
            profile = await analyzer.profile_file(file_path, "xlsx")
            parts.append(analyzer.format_profile(profile))
            return "\n".join(parts)
        except ImportError:
            pass
        except Exception:
            pass

        # Fallback to xlrd for .xls
        try:
            import xlrd
            wb = xlrd.open_workbook(file_path)
            parts = []
            for sheet in wb.sheets():
                parts.append(f"\n[工作表: {sheet.name}]")
                for row_idx in range(min(sheet.nrows, 500)):
                    row = sheet.row_values(row_idx)
                    row_text = " | ".join(str(c) for c in row)
                    if row_text.strip():
                        parts.append(row_text)
            profile = await analyzer.profile_file(file_path, "xls")
            parts.append(analyzer.format_profile(profile))
            return "\n".join(parts)
        except ImportError:
            return "[openpyxl/xlrd未安装，无法提取Excel内容]"
        except Exception as e:
            return f"[Excel提取失败: {str(e)}]"

    async def _extract_csv(self, file_path: str) -> str:
        """Extract text from CSV file."""
        try:
            # Detect encoding
            encoding = await self._detect_encoding(file_path)
            parts = []
            row_count = 0
            with open(file_path, "r", encoding=encoding, errors="replace") as f:
                reader = csv.reader(f)
                for row in reader:
                    parts.append(" | ".join(row))
                    row_count += 1
                    if row_count > 1000:
                        parts.append(f"[... 仅显示前1000行 ...]")
                        break
            analyzer = ExcelAnalyzer()
            profile = await analyzer.profile_file(file_path, "csv")
            parts.append(analyzer.format_profile(profile))
            return "\n".join(parts)
        except Exception as e:
            return f"[CSV提取失败: {str(e)}]"

    async def _extract_pptx(self, file_path: str) -> str:
        """Extract text from PPTX using python-pptx."""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            parts = []
            for i, slide in enumerate(prs.slides):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                if slide_texts:
                    parts.append(f"\n[幻灯片 {i+1}]\n" + "\n".join(slide_texts))
            return "\n".join(parts)
        except ImportError:
            return "[python-pptx未安装，无法提取PPTX内容]"
        except Exception as e:
            return f"[PPTX提取失败: {str(e)}]"

    async def _extract_json(self, file_path: str) -> str:
        """Extract text from JSON file."""
        try:
            encoding = await self._detect_encoding(file_path)
            with open(file_path, "r", encoding=encoding, errors="replace") as f:
                data = json.load(f)
            return json.dumps(data, ensure_ascii=False, indent=2)
        except json.JSONDecodeError:
            return await self._extract_text_file(file_path)
        except Exception as e:
            return f"[JSON提取失败: {str(e)}]"

    async def _extract_xml(self, file_path: str) -> str:
        """Extract text from XML/HTML file."""
        try:
            from lxml import etree
            with open(file_path, "rb") as f:
                content = f.read()
            try:
                root = etree.fromstring(content)
                # Extract all text nodes
                texts = root.xpath("//text()")
                return " ".join(t.strip() for t in texts if t.strip())
            except Exception:
                # Try as HTML
                from lxml import html
                tree = html.fromstring(content)
                return tree.text_content()
        except ImportError:
            return await self._extract_text_file(file_path)
        except Exception as e:
            return f"[XML/HTML提取失败: {str(e)}]"

    async def _extract_text_file(self, file_path: str) -> str:
        """Read a plain text file."""
        encoding = await self._detect_encoding(file_path)
        with open(file_path, "r", encoding=encoding, errors="replace") as f:
            return f.read()

    async def _detect_encoding(self, file_path: str) -> str:
        """Detect file encoding using chardet."""
        try:
            import chardet
            with open(file_path, "rb") as f:
                raw = f.read(10000)
            result = chardet.detect(raw)
            encoding = result.get("encoding") or "utf-8"
            # Map some common encodings
            if encoding.lower() in ("ascii",):
                encoding = "utf-8"
            return encoding
        except ImportError:
            return "utf-8"
        except Exception:
            return "utf-8"

    def get_supported_extensions(self) -> list:
        """Return list of supported file extensions."""
        return [
            ".pdf", ".docx", ".doc", ".xlsx", ".xls", ".csv",
            ".pptx", ".json", ".xml", ".html", ".htm",
            ".txt", ".md", ".py", ".js", ".ts", ".java",
            ".cpp", ".c", ".go", ".rs", ".sql", ".sh",
            ".yaml", ".yml", ".toml", ".ini", ".cfg",
        ]
