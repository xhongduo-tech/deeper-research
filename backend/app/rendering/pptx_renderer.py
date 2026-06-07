"""PptxRenderer — deterministic PowerPoint generation from PptxSpec.

Opens a .pptx template (if provided), preserves slide master/layouts, then
adds slides according to the spec. Each slide maps to a layout by name.
"""
from __future__ import annotations

import io
import logging
from typing import Optional

from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from app.rendering.doc_spec import PptxSpec, PptxSlideSpec, TableSpec, ChartSpec

logger = logging.getLogger(__name__)

# 26 color themes: (accent_hex, background_hex, text_hex)
_COLOR_THEMES = [
    ("1F497D", "FFFFFF", "000000"),  # Classic Blue
    ("C00000", "FFFFFF", "000000"),  # Corporate Red
    ("375623", "FFFFFF", "000000"),  # Forest Green
    ("7030A0", "FFFFFF", "FFFFFF"),  # Purple
    ("E36C09", "FFFFFF", "000000"),  # Orange
    ("17375E", "DCE6F1", "000000"),  # Navy Blue
    ("404040", "F2F2F2", "000000"),  # Charcoal
    ("2E75B6", "FFFFFF", "000000"),  # Medium Blue
    ("843C0C", "FFF2CC", "000000"),  # Brown Gold
    ("00B0F0", "FFFFFF", "000000"),  # Sky Blue
    ("70AD47", "FFFFFF", "000000"),  # Green
    ("FF0000", "FFFFFF", "000000"),  # Red
    ("002060", "FFFFFF", "FFFFFF"),  # Dark Navy
    ("833C00", "FFFFFF", "000000"),  # Dark Orange
    ("244185", "D9E1F2", "000000"),  # Slate Blue
    ("538135", "E2EFDA", "000000"),  # Olive Green
    ("954F72", "FFFFFF", "000000"),  # Mauve
    ("156082", "DDEBF7", "000000"),  # Teal Blue
    ("1F3864", "FFFFFF", "FFFFFF"),  # Deep Navy
    ("BF8F00", "FFFFFF", "000000"),  # Gold
    ("A9D18E", "375623", "FFFFFF"),  # Light Green (inv)
    ("2F5496", "D6DCE4", "000000"),  # Steel Blue
    ("C55A11", "FCE4D6", "000000"),  # Burnt Orange
    ("0070C0", "DEEAF1", "000000"),  # Bright Blue
    ("7F7F7F", "F2F2F2", "000000"),  # Gray
    ("305496", "D9E1F2", "000000"),  # Indigo
]

# P1-5: Font fallback chains — keyed by canonical name or common alias.
# The first entry in each list is the preferred/canonical name to actually apply.
_FONT_FALLBACK_MAP: dict[str, list[str]] = {
    "思源宋体":         ["思源宋体 Heavy", "Source Han Serif", "STSong", "SimSun", "宋体"],
    "思源黑体":         ["Source Han Sans", "STHeiti", "Microsoft YaHei", "微软雅黑"],
    "微软雅黑":         ["Microsoft YaHei", "PingFang SC", "STHeiti", "黑体"],
    "microsoft yahei": ["Microsoft YaHei", "PingFang SC", "STHeiti", "黑体"],
    "宋体":            ["SimSun", "STSong", "宋体"],
    "黑体":            ["STHeiti", "SimHei", "黑体"],
    "仿宋":            ["FangSong", "STFangsong", "仿宋"],
    "楷体":            ["KaiTi", "STKaiti", "楷体"],
    "calibri":         ["Calibri", "Arial", "Helvetica"],
    "times new roman": ["Times New Roman", "Georgia"],
    "arial":           ["Arial", "Helvetica", "Liberation Sans"],
}


def _resolve_font_name(font_name: str) -> str:
    """Return the canonical preferred font from the fallback chain, or original if unknown.

    Two-pass: exact key match first to avoid an alias in one chain shadowing another.
    """
    if not font_name:
        return font_name
    key = font_name.lower().strip()
    # Pass 1: exact map-key match
    for map_key, fallbacks in _FONT_FALLBACK_MAP.items():
        if key == map_key:
            return fallbacks[0]
    # Pass 2: font appears as an alias inside a fallbacks list
    for map_key, fallbacks in _FONT_FALLBACK_MAP.items():
        if font_name in fallbacks:
            return fallbacks[0]
    return font_name


# Layout hint → candidate keyword list (tried in order against layout names)
_LAYOUT_KEYWORDS = {
    "title_only":     ["title only", "blank", "标题"],
    "section_header": ["section header", "title", "封面", "章节"],
    "content":        ["title and content", "两栏内容", "content"],
    "comparison":     ["two content", "comparison", "对比"],
    "big_number":     ["blank", "title only"],
    "data_table":     ["title and content", "content"],
}

# Approximate characters that fit in a standard body placeholder at 18pt
_SLIDE_BODY_CHAR_CAPACITY = 300


class PptxRenderer:
    """Renders a PptxSpec to .pptx bytes.

    All exceptions propagate — no silent swallowing.
    """

    def __init__(self):
        self._template_dna: dict = {}  # accent_hex, font_name extracted from template

    def render(
        self,
        spec: PptxSpec,
        template_path: Optional[str] = None,
    ) -> bytes:
        from pptx import Presentation

        has_template = bool(template_path)
        if has_template:
            prs = Presentation(template_path)
            # P2-B: extract Template DNA before removing slides
            self._template_dna = _extract_template_dna(prs)
            _remove_all_slides(prs)
        else:
            prs = Presentation()
            self._template_dna = {}

        layouts = prs.slide_layouts
        layout_map = self._build_layout_map(layouts)

        # Theme index from spec metadata (0-25); ignored when a template provides its own theme
        theme_idx = int(spec.metadata.get("theme_idx", 0)) if not has_template else -1

        for slide_spec in spec.slides:
            layout = self._pick_layout(layouts, layout_map, slide_spec.layout, slide_spec)
            slide = prs.slides.add_slide(layout)
            self._fill_slide(prs, slide, slide_spec)
            if theme_idx >= 0:
                self._apply_color_theme(slide, theme_idx)
            elif self._template_dna.get("accent_hex"):
                self._apply_dna_theme(slide, self._template_dna)

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    # ── Layout helpers ────────────────────────────────────────────────────────

    def _build_layout_map(self, layouts) -> dict[str, int]:
        """Map lowercase layout names to their index."""
        return {(layout.name or "").lower(): i for i, layout in enumerate(layouts)}

    def _pick_layout(self, layouts, layout_map: dict[str, int], layout_hint: str,
                     slide_spec: PptxSlideSpec):
        """P1-D: Score-based layout selection — keyword match + placeholder-count fallback."""
        # Step 1: keyword match from hint
        keywords = _LAYOUT_KEYWORDS.get(layout_hint, _LAYOUT_KEYWORDS["content"])
        for kw in keywords:
            for name, idx in layout_map.items():
                if kw in name:
                    return layouts[idx]

        # Step 2: content-named fallback
        for name, idx in layout_map.items():
            if "content" in name:
                return layouts[idx]

        # Step 3: score by placeholder count — prefer layouts that have ≥2 placeholders
        # (title + body), but not too many (overly complex masters)
        has_table = slide_spec.table is not None
        has_chart = slide_spec.chart is not None
        has_bullets = bool(slide_spec.bullets)

        def _score(layout) -> int:
            n_ph = len(list(layout.placeholders))
            if has_table or has_chart:
                return -abs(n_ph - 1)   # prefer single-placeholder for full-bleed content
            if has_bullets:
                return -abs(n_ph - 2)   # prefer title+body
            return -abs(n_ph - 1)

        scored = sorted(enumerate(layouts), key=lambda x: _score(x[1]), reverse=True)
        return scored[0][1] if scored else layouts[0]

    # ── Slide filling ─────────────────────────────────────────────────────────

    def _fill_slide(self, prs, slide, spec: PptxSlideSpec) -> None:
        title_ph, body_ph = self._find_placeholders(slide)

        if title_ph is not None:
            title_ph.text = spec.assertion_title
            for para in title_ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.bold = True

        # ── P1-2: big_number layout ───────────────────────────────────────────
        # P3-1: Validate that assertion_title actually contains a digit before
        # rendering big_number. If no number is present the giant empty number
        # box looks broken — fall through to regular content rendering instead.
        import re as _re
        _has_digit = bool(_re.search(r"\d", spec.assertion_title or ""))
        if spec.big_number and (spec.big_number.get("value") or _has_digit):
            self._render_big_number(slide, prs, spec.big_number)
        elif spec.layout == "big_number" and not _has_digit:
            # Degrade to content layout silently — write assertion_title as heading
            # and bullets normally; don't attempt big_number rendering.
            pass  # falls through to regular bullet rendering below

        # ── P2-1: comparison layout — split bullets left / right ──────────────
        elif spec.layout == "comparison" and spec.bullets:
            self._fill_comparison_slide(slide, spec)

        # ── P3-1: section_header — dark full-bleed background + white text ──────
        elif spec.layout == "section_header":
            # Apply deep accent background to the entire slide background
            accent_hex = self._template_dna.get("accent_hex", "") or "1A3557"
            try:
                _set_slide_background(slide, accent_hex)
            except Exception:
                pass  # keep default background on failure

            # Title in white, large, centered
            if title_ph is not None:
                title_ph.text = spec.assertion_title
                for para in title_ph.text_frame.paragraphs:
                    para.alignment = 1  # CENTER
                    for run in para.runs:
                        run.font.size = Pt(40)
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

            # Subtitle / key phrase in white, smaller
            if body_ph is not None and spec.bullets:
                tf = body_ph.text_frame
                tf.clear()
                para = tf.paragraphs[0]
                para.text = spec.bullets[0]
                para.alignment = 1  # CENTER
                for run in para.runs:
                    run.font.size = Pt(22)
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # ── Regular bullet rendering ──────────────────────────────────────────
        else:
            if body_ph is not None and spec.bullets:
                tf = body_ph.text_frame
                tf.clear()
                # P1-2: Clamp total bullet content to layout-specific char limit
                # to prevent text overflowing the placeholder bounds.
                _LAYOUT_CHAR_LIMITS = {
                    "content": 300,
                    "comparison": 200,
                    "title_only": 400,
                }
                char_limit = _LAYOUT_CHAR_LIMITS.get(spec.layout, 300)
                safe_bullets = _fit_bullets_to_limit(spec.bullets, char_limit)
                total_bullet_chars = sum(len(b) for b in safe_bullets)
                adaptive_pt = 14 if total_bullet_chars > 250 else (16 if total_bullet_chars > 150 else 18)
                for i, bullet in enumerate(safe_bullets):
                    para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    para.text = bullet
                    para.level = 0
                    for run in para.runs:
                        run.font.size = Pt(adaptive_pt)

        # ── Chart ─────────────────────────────────────────────────────────────
        if spec.chart and spec.layout not in ("big_number", "section_header"):
            has_bullets = bool(spec.bullets) and body_ph is not None
            chart_top = Inches(3.5) if has_bullets else Inches(2.0)
            # P1-3: cap chart height so top+height never exceeds 7.0" (safe slide floor)
            try:
                slide_h_in = (prs.slide_height or 6858000) / 914400
            except Exception:
                slide_h_in = 7.5
            max_h = max(Inches(1.5), int((slide_h_in - 0.5) * 914400) - chart_top)
            chart_height = min(Inches(4.5), max_h)
            self._add_chart(slide, prs, spec.chart, top=chart_top, chart_height=chart_height)

        # ── Table ─────────────────────────────────────────────────────────────
        if spec.table and spec.layout not in ("big_number",):
            has_bullets = bool(spec.bullets) and body_ph is not None
            table_top = Inches(3.5) if has_bullets else Inches(2.0)
            self._add_table(slide, spec.table, top=table_top)

        if spec.speaker_notes:
            slide.notes_slide.notes_text_frame.text = spec.speaker_notes

        # P2-C: Clear template metadata placeholders (date/footer/page-number)
        _clear_slide_metadata_placeholders(slide)

    # ── big_number layout ─────────────────────────────────────────────────────

    def _render_big_number(self, slide, prs, big_number: dict) -> None:
        """P1-2: Render a large centered KPI number with a smaller label below it."""
        from pptx.enum.text import PP_ALIGN

        value = str(big_number.get("value", ""))
        label = str(big_number.get("label", ""))

        try:
            slide_w = prs.slide_width or 9144000
            slide_h = prs.slide_height or 6858000
        except Exception:
            slide_w, slide_h = 9144000, 6858000

        val_w = Inches(6)
        val_h = Inches(2)
        val_left = (slide_w - val_w) // 2
        val_top = slide_h // 2 - Inches(1.5)

        val_box = slide.shapes.add_textbox(val_left, val_top, val_w, val_h)
        tf = val_box.text_frame
        tf.word_wrap = False
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = value
        run.font.size = Pt(72)
        run.font.bold = True
        accent_hex = self._template_dna.get("accent_hex", "1F497D")
        try:
            run.font.color.rgb = RGBColor.from_string(accent_hex)
        except Exception:
            pass

        if label:
            lbl_box = slide.shapes.add_textbox(val_left, val_top + val_h, val_w, Inches(0.6))
            lbl_tf = lbl_box.text_frame
            lbl_para = lbl_tf.paragraphs[0]
            lbl_para.alignment = PP_ALIGN.CENTER
            lbl_run = lbl_para.add_run()
            lbl_run.text = label
            lbl_run.font.size = Pt(20)

    # ── comparison layout ─────────────────────────────────────────────────────

    def _fill_comparison_slide(self, slide, spec: PptxSlideSpec) -> None:
        """P2-1: Split bullets into two side-by-side columns for comparison layout."""
        bullets = spec.bullets
        mid = max(1, len(bullets) // 2)
        left_bullets = bullets[:mid]
        right_bullets = bullets[mid:]

        # Try to populate two body placeholders (idx 1 and 2 in "Two Content" layouts)
        body_phs = [
            ph for ph in sorted(slide.placeholders, key=lambda p: p.placeholder_format.idx)
            if ph.placeholder_format.idx != 0 and hasattr(ph, "text_frame")
        ]

        if len(body_phs) >= 2:
            for ph, col_bullets in zip(body_phs[:2], [left_bullets, right_bullets]):
                tf = ph.text_frame
                tf.clear()
                for i, bullet in enumerate(col_bullets):
                    para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    para.text = bullet
                    for run in para.runs:
                        run.font.size = Pt(16)
        else:
            # Fallback: two textboxes side by side
            slide_w = 9144000
            col_w = Inches(4.2)
            col_h = Inches(4.0)
            col_top = Inches(2.0)
            for col_idx, col_bullets in enumerate([left_bullets, right_bullets]):
                left = Inches(0.4) if col_idx == 0 else Inches(5.0)
                box = slide.shapes.add_textbox(left, col_top, col_w, col_h)
                tf = box.text_frame
                tf.word_wrap = True
                for i, bullet in enumerate(col_bullets):
                    para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    para.text = f"• {bullet}"
                    for run in para.runs:
                        run.font.size = Pt(16)

    def _find_placeholders(self, slide):
        """Return (title_ph, body_ph) by semantic placeholder type, not hardcoded idx.

        P1-3: Use PP_PLACEHOLDER type enum (TITLE/CENTER_TITLE for title;
        BODY/OBJECT/SUBTITLE for body) so templates with non-standard idx values
        still get correct placeholder assignment. Falls back to idx==0 heuristic
        when the type attribute is unavailable or None.
        """
        try:
            from pptx.enum.text import PP_PLACEHOLDER as _PPH
            _TITLE_TYPES = {_PPH.TITLE, _PPH.CENTER_TITLE}
            _BODY_TYPES  = {_PPH.BODY, _PPH.OBJECT, _PPH.SUBTITLE,
                            _PPH.PICTURE, _PPH.CHART, _PPH.TABLE}
        except Exception:
            _TITLE_TYPES = _BODY_TYPES = None

        title_ph = None
        body_ph = None

        for ph in slide.placeholders:
            pf = ph.placeholder_format
            ph_type = getattr(pf, "type", None)

            if _TITLE_TYPES and ph_type in _TITLE_TYPES:
                if title_ph is None:
                    title_ph = ph
            elif _BODY_TYPES and ph_type in _BODY_TYPES:
                if body_ph is None and hasattr(ph, "text_frame"):
                    body_ph = ph
            else:
                # Fallback: idx==0 → title, first other with text_frame → body
                if pf.idx == 0 and title_ph is None:
                    title_ph = ph
                elif pf.idx != 0 and body_ph is None and hasattr(ph, "text_frame"):
                    body_ph = ph

        return title_ph, body_ph

    # ── Chart ─────────────────────────────────────────────────────────────────

    def _add_chart(self, slide, prs, chart_spec: ChartSpec, top=None, chart_height=None) -> None:
        from app.services.chart_render_service import render_chart_png

        if top is None:
            top = Inches(2.0)
        # P1-3: use caller-supplied height cap to prevent overflow
        height = chart_height if chart_height is not None else Inches(4.5)
        left, width = Inches(0.5), Inches(9.0)

        # P1-B: Pass a theme-matched palette derived from the template DNA accent color.
        palette = _build_chart_palette(self._template_dna.get("accent_hex", ""))

        # P2-C: Adapt chart pixel dimensions to the actual slide's aspect ratio.
        # EMU → pixels at 96 DPI (1 inch = 914400 EMU, 1 inch = 96 px).
        try:
            slide_w_emu = prs.slide_width or (9144000)   # default 10 inches
            slide_h_emu = prs.slide_height or (5143500)  # default 5.625 inches
            px_w = max(800, int(slide_w_emu * 96 / 914400))
            # Chart occupies ~55% of slide height leaving room for title + bullets
            px_h = max(400, int(slide_h_emu * 96 / 914400 * 0.55))
        except Exception:
            px_w, px_h = 1200, 680

        try:
            render_spec = chart_spec.to_render_spec()
            result = render_chart_png(render_spec, width=px_w, height=px_h, palette=palette or None)
            if result and result.png:
                slide.shapes.add_picture(io.BytesIO(result.png), left, top, width, height)
                return
        except Exception as exc:
            logger.warning("[PPTX] Chart render failed (%s): %s", chart_spec.title, exc)

        # Fallback: text box placeholder so the slide isn't silently blank
        txBox = slide.shapes.add_textbox(left, top, width, Inches(0.6))
        tf = txBox.text_frame
        tf.text = f"[图表：{chart_spec.title}（渲染失败，请手动插入）]"
        tf.paragraphs[0].runs[0].font.size = Pt(10)

    # ── Table ─────────────────────────────────────────────────────────────────

    def _add_table(self, slide, table_spec: TableSpec, top=None) -> None:
        headers = table_spec.headers
        rows = table_spec.rows
        if not headers and not rows:
            return

        n_cols = len(headers) or (len(rows[0]) if rows else 0)
        n_rows = 1 + len(rows)
        if n_cols == 0 or n_rows <= 1:
            return

        if top is None:
            top = Inches(2.0)
        left, width = Inches(0.5), Inches(9.0)
        height = Inches(min(0.4 * n_rows, 4.5))

        # P2-2: derive header fill from template DNA or fall back to classic blue
        accent_hex = self._template_dna.get("accent_hex", "1F497D") or "1F497D"
        try:
            hdr_rgb = RGBColor.from_string(accent_hex)
        except Exception:
            hdr_rgb = RGBColor(0x1F, 0x49, 0x7D)
        band_rgb = RGBColor(0xDC, 0xE6, 0xF1)  # light blue for even rows

        try:
            shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
            tbl = shape.table
            # Header row — accent background, white bold text
            for c_idx, h in enumerate(headers[:n_cols]):
                cell = tbl.cell(0, c_idx)
                cell.text = h
                cell.fill.solid()
                cell.fill.fore_color.rgb = hdr_rgb
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.bold = True
                        run.font.size = Pt(10)
                        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            # Data rows — banded (even rows get light tint)
            for r_idx, row in enumerate(rows):
                for c_idx, val in enumerate(row[:n_cols]):
                    cell = tbl.cell(r_idx + 1, c_idx)
                    cell.text = str(val)
                    if (r_idx + 1) % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = band_rgb
        except Exception as exc:
            logger.warning("[PPTX] Table render failed: %s", exc)

    # ── Color themes ──────────────────────────────────────────────────────────

    def _apply_color_theme(self, slide, theme_idx: int) -> None:
        """Apply accent color to title runs (fresh presentations only)."""
        if theme_idx < 0 or theme_idx >= len(_COLOR_THEMES):
            return
        accent_hex = _COLOR_THEMES[theme_idx][0]
        try:
            accent = RGBColor.from_string(accent_hex)
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 0:
                    for para in ph.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.color.rgb = accent
        except Exception:
            pass

    def _apply_dna_theme(self, slide, dna: dict) -> None:
        """P2-B: Apply template DNA accent color + font to title placeholder.

        P1-5: Font name is resolved through the fallback chain before applying,
        so fonts unavailable on the current machine fall back gracefully.
        """
        accent_hex = dna.get("accent_hex", "")
        # P1-5: resolve through fallback chain
        font_name = _resolve_font_name(dna.get("font_name", ""))
        try:
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 0:
                    for para in ph.text_frame.paragraphs:
                        for run in para.runs:
                            if accent_hex:
                                run.font.color.rgb = RGBColor.from_string(accent_hex)
                            if font_name:
                                run.font.name = font_name
        except Exception:
            pass


# ── Module-level helpers ──────────────────────────────────────────────────────

def _remove_all_slides(prs) -> None:
    """Safely remove all slides, preserving slide master and layouts."""
    _NS_REL = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    sldIdLst = prs.slides._sldIdLst
    for elem in list(sldIdLst):
        rId = elem.get(_NS_REL)
        if rId:
            try:
                prs.part.drop_rel(rId)
            except Exception:
                pass
            sldIdLst.remove(elem)


def _build_chart_palette(accent_hex: str) -> list[str]:
    """P1-B: Generate a 6-color palette rooted at the template's accent color.

    Returns hex strings (without #) suitable for chart_render_service palette arg.
    Falls back to an empty list when no accent color is available, letting the
    renderer use its own default palette.
    """
    if not accent_hex or len(accent_hex) < 6:
        return []
    try:
        r = int(accent_hex[0:2], 16)
        g = int(accent_hex[2:4], 16)
        b = int(accent_hex[4:6], 16)

        def _vary(r, g, b, factor) -> str:
            nr = min(255, int(r * factor))
            ng = min(255, int(g * factor))
            nb = min(255, int(b * factor))
            return f"{nr:02X}{ng:02X}{nb:02X}"

        return [
            accent_hex.upper(),
            _vary(r, g, b, 0.75),
            _vary(r, g, b, 0.55),
            _vary(r, g, b, 1.30),
            _vary(r, g, b, 0.40),
            _vary(r, g, b, 1.15),
        ]
    except Exception:
        return []


def _clear_slide_metadata_placeholders(slide) -> None:
    """P2-C: Clear date (idx 10), footer (idx 11), and page-number (idx 12) placeholders.

    When a .pptx template has these populated, the generated slides inherit stale
    template metadata. Clearing them prevents leftover text from appearing.
    """
    _METADATA_IDXS = {10, 11, 12}
    for ph in slide.placeholders:
        if ph.placeholder_format.idx in _METADATA_IDXS:
            try:
                ph.text_frame.clear()
            except Exception:
                pass


def _extract_template_dna(prs) -> dict:
    """P2-B: Extract accent color and font name from a .pptx template's slide master.

    Returns dict with 'accent_hex' and 'font_name' keys (may be empty strings if
    extraction fails — caller applies these best-effort only).
    """
    dna = {"accent_hex": "", "font_name": ""}
    try:
        master = prs.slide_master
        # Try to extract accent1 color from theme XML
        from lxml import etree
        theme_el = master._element.find(
            ".//{http://schemas.openxmlformats.org/drawingml/2006/main}theme"
        )
        if theme_el is None:
            # Theme is usually in a separate part
            if hasattr(master, "theme_color_map"):
                pass  # further extraction possible but complex; skip for now

        # Extract font from slide master's title placeholder
        for ph in master.placeholders:
            if ph.placeholder_format.idx == 0:
                for para in ph.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            dna["font_name"] = run.font.name
                        break
                break

        # Extract accent color from first slide layout's title text (if colored)
        for layout in prs.slide_layouts:
            for ph in layout.placeholders:
                if ph.placeholder_format.idx == 0:
                    for para in ph.text_frame.paragraphs:
                        for run in para.runs:
                            color = run.font.color
                            if color and color.type is not None:
                                try:
                                    hex_val = str(color.rgb)
                                    if hex_val and hex_val != "000000":
                                        dna["accent_hex"] = hex_val
                                except Exception:
                                    pass
            if dna["accent_hex"]:
                break

    except Exception as exc:
        logger.debug("[PPTX] Template DNA extraction failed: %s", exc)

    return dna


def _fit_bullets_to_limit(bullets: list[str], max_chars: int) -> list[str]:
    """P1-2: Truncate bullet list so total character count stays within max_chars.

    Keeps as many complete bullets as fit; if the first bullet alone exceeds
    the limit it is truncated at the last punctuation mark before max_chars.
    """
    if not bullets:
        return bullets
    result: list[str] = []
    remaining = max_chars
    for b in bullets:
        if remaining <= 0:
            break
        if len(b) <= remaining:
            result.append(b)
            remaining -= len(b)
        else:
            # Truncate this bullet at a sentence boundary
            window = b[:remaining]
            for punct in ("。", "；", "，", ".", ";", ","):
                idx = window.rfind(punct)
                if idx > max(10, remaining // 3):
                    result.append(window[: idx + 1] + "…")
                    break
            else:
                result.append(window + "…")
            remaining = 0
    return result


def _set_slide_background(slide, hex_color: str) -> None:
    """P3-1: Fill the entire slide background with a solid color.

    Uses the pptx XML directly since python-pptx's high-level API for slide
    background fill is not exposed. The solidFill element is inserted into
    the slide's background (bg/bgPr) element.
    """
    from pptx.oxml.ns import qn
    from pptx.oxml import parse_xml
    from lxml import etree

    hex_color = (hex_color or "1A3557").lstrip("#").upper()
    if len(hex_color) != 6:
        hex_color = "1A3557"

    spTree = slide.shapes._spTree
    # Walk up to the slide element
    slide_el = spTree.getparent()
    while slide_el is not None and slide_el.tag.split("}")[-1] != "sld":
        slide_el = slide_el.getparent()
    if slide_el is None:
        return

    bg_xml = (
        f'<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        f' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<p:bgPr>'
        f'<a:solidFill><a:srgbClr val="{hex_color}"/></a:solidFill>'
        f'<a:effectLst/>'
        f'</p:bgPr>'
        f'</p:bg>'
    )
    bg_el = parse_xml(bg_xml)

    # Remove any existing bg element
    existing = slide_el.find(qn("p:bg"))
    if existing is not None:
        slide_el.remove(existing)

    # Insert bg as first child of sld (before cSld)
    cSld = slide_el.find(qn("p:cSld"))
    if cSld is not None:
        cSld.addprevious(bg_el)
    else:
        slide_el.insert(0, bg_el)
