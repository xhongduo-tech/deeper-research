"""Deterministic reference document structure extraction — no LLM involved.

Extracts headings, section count, style characteristics from uploaded .docx/.pptx
files so the PLAN phase can derive the outline directly from the template structure
rather than having the LLM re-invent it.
"""
from __future__ import annotations

import re
from dataclasses import dataclass, field
from typing import Optional


# ── Data model ────────────────────────────────────────────────────────────────

@dataclass
class ReferenceSection:
    heading_text: str       # exact heading text from the original document
    heading_level: int      # 1 = H1, 2 = H2
    content_type: str       # "paragraphs" | "table" | "mixed"
    char_count: int         # rough content length in characters
    has_table: bool
    has_bullets: bool
    paragraph_count: int


@dataclass
class ReferenceStructure:
    doc_type: str                           # "speech" | "report" | "slide_deck"
    sections: list[ReferenceSection] = field(default_factory=list)
    style_note: str = "structured_headings" # "narrative_prose" | "structured_headings"
    estimated_total_chars: int = 0
    slide_count: Optional[int] = None       # PPTX only
    # P2-5: Typography extracted from template for style replication
    body_font: str = ""                     # Normal paragraph font name
    heading_font: str = ""                  # Heading 1 font name
    body_size_pt: float = 0.0              # Normal font size in pt
    body_space_before_pt: float = 0.0      # Paragraph space-before in pt
    body_space_after_pt: float = 0.0       # Paragraph space-after in pt

    def to_dict(self) -> dict:
        """Serialise to a plain dict for storage in report.scoping_plan."""
        return {
            "doc_type": self.doc_type,
            "style_note": self.style_note,
            "estimated_total_chars": self.estimated_total_chars,
            "slide_count": self.slide_count,
            "body_font": self.body_font,
            "heading_font": self.heading_font,
            "body_size_pt": self.body_size_pt,
            "body_space_before_pt": self.body_space_before_pt,
            "body_space_after_pt": self.body_space_after_pt,
            "sections": [
                {
                    "heading_text": s.heading_text,
                    "heading_level": s.heading_level,
                    "content_type": s.content_type,
                    "char_count": s.char_count,
                    "has_table": s.has_table,
                    "has_bullets": s.has_bullets,
                    "paragraph_count": s.paragraph_count,
                }
                for s in self.sections
            ],
        }

    @classmethod
    def from_dict(cls, d: dict) -> "ReferenceStructure":
        """Deserialise from report.scoping_plan storage."""
        sections = [
            ReferenceSection(**s) for s in d.get("sections", [])
        ]
        return cls(
            doc_type=d.get("doc_type", "report"),
            sections=sections,
            style_note=d.get("style_note", "structured_headings"),
            estimated_total_chars=d.get("estimated_total_chars", 0),
            slide_count=d.get("slide_count"),
            body_font=d.get("body_font", ""),
            heading_font=d.get("heading_font", ""),
            body_size_pt=d.get("body_size_pt", 0.0),
            body_space_before_pt=d.get("body_space_before_pt", 0.0),
            body_space_after_pt=d.get("body_space_after_pt", 0.0),
        )


# ── DOCX extraction ──────────────────────────────────────────────────────────

# Matches Chinese-style numbered headings like "一、", "二、", "(一)", "1."
_CN_HEADING_RE = re.compile(
    r"^(?:[一二三四五六七八九十百千]{1,3}[、。]"
    r"|第[一二三四五六七八九十百千]{1,3}[章节部分]"
    r"|（[一二三四五六七八九十]{1,2}）"
    r"|\d{1,2}[、.．]\s?"
    r"|\([一二三四五六七八九十\d]{1,2}\))"
)

# Bullet indicators
_BULLET_RE = re.compile(r"^[•·◆◇▪▸➤\-\*]\s|^[①②③④⑤⑥⑦⑧⑨⑩]")


def extract_docx_structure(file_path: str) -> Optional[ReferenceStructure]:
    """Extract heading / section structure from a .docx file.

    Returns None if python-docx is unavailable or the file cannot be parsed.
    """
    try:
        from docx import Document
        from docx.oxml.ns import qn
    except ImportError:
        return None

    try:
        doc = Document(file_path)
    except Exception:
        return None

    # Collect headings and their following content paragraphs
    sections: list[ReferenceSection] = []
    current_heading: Optional[str] = None
    current_level: int = 1
    current_paras: list[str] = []
    has_table_in_section: bool = False

    # Build a set of paragraph positions that are inside tables (to skip)
    table_para_set: set[int] = set()
    for tbl in doc.tables:
        for row in tbl.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    table_para_set.add(id(para))

    def _flush_section():
        nonlocal current_heading, current_paras, has_table_in_section
        if current_heading is None:
            return
        text_blob = " ".join(current_paras)
        has_bullets = any(_BULLET_RE.match(p) for p in current_paras)
        content_type = (
            "table" if has_table_in_section and not current_paras
            else "mixed" if has_table_in_section
            else "paragraphs"
        )
        sections.append(ReferenceSection(
            heading_text=current_heading,
            heading_level=current_level,
            content_type=content_type,
            char_count=len(text_blob),
            has_table=has_table_in_section,
            has_bullets=has_bullets,
            paragraph_count=len(current_paras),
        ))
        current_heading = None
        current_paras = []
        has_table_in_section = False

    # Track which tables we encounter between headings
    para_list = list(doc.paragraphs)
    # Also note table positions by interleaving with XML order
    body = doc.element.body
    table_ids_in_order: list[int] = []
    for child in body:
        if child.tag == qn("w:tbl"):
            table_ids_in_order.append(id(child))

    seen_table_ids: set[int] = set()

    for para in para_list:
        if id(para) in table_para_set:
            continue  # skip table cells

        style_name = (para.style.name or "").lower()
        text = para.text.strip()
        if not text:
            continue

        is_heading = (
            "heading" in style_name
            or style_name.startswith("标题")
            or style_name in ("title", "subtitle")
        )
        # Also treat short all-bold lines or CN-numbered lines as headings
        if not is_heading and text:
            is_cn_heading = bool(_CN_HEADING_RE.match(text)) and len(text) <= 40
            all_bold = all(run.bold for run in para.runs if run.text.strip())
            if is_cn_heading or (all_bold and len(text) <= 30 and para.runs):
                is_heading = True

        if is_heading:
            _flush_section()
            current_heading = text
            # Infer level from style name number or heuristic
            m = re.search(r"\d", style_name)
            current_level = int(m.group()) if m else 1
        else:
            if current_heading is not None:
                current_paras.append(text)

    _flush_section()

    if not sections:
        # Speechs / narrative prose: no headings — treat the whole doc as one section
        all_text = " ".join(p.text.strip() for p in doc.paragraphs if p.text.strip())
        sections.append(ReferenceSection(
            heading_text="（全文）",
            heading_level=1,
            content_type="paragraphs",
            char_count=len(all_text),
            has_table=bool(doc.tables),
            has_bullets=False,
            paragraph_count=len([p for p in doc.paragraphs if p.text.strip()]),
        ))
        style_note = "narrative_prose"
        doc_type = "speech"
    else:
        style_note = "structured_headings"
        doc_type = "report"

    total_chars = sum(s.char_count for s in sections)

    # P2-5: Extract typography from Normal / Heading 1 styles
    body_font = body_size_pt = body_space_before_pt = body_space_after_pt = heading_font = ""
    body_size_pt = body_space_before_pt = body_space_after_pt = 0.0
    try:
        normal = doc.styles["Normal"]
        if normal.font.name:
            body_font = normal.font.name
        if normal.font.size:
            body_size_pt = normal.font.size.pt
        pf = normal.paragraph_format
        if pf.space_before:
            body_space_before_pt = pf.space_before.pt
        if pf.space_after:
            body_space_after_pt = pf.space_after.pt
    except Exception:
        pass
    try:
        h1 = doc.styles["Heading 1"]
        if h1.font.name:
            heading_font = h1.font.name
        elif h1.base_style and h1.base_style.font.name:
            heading_font = h1.base_style.font.name
    except Exception:
        pass

    return ReferenceStructure(
        doc_type=doc_type,
        sections=sections,
        style_note=style_note,
        estimated_total_chars=total_chars,
        body_font=body_font,
        heading_font=heading_font,
        body_size_pt=body_size_pt,
        body_space_before_pt=body_space_before_pt,
        body_space_after_pt=body_space_after_pt,
    )


# ── PPTX extraction ──────────────────────────────────────────────────────────

def extract_pptx_structure(file_path: str) -> Optional[ReferenceStructure]:
    """Extract slide structure from a .pptx file."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        return None

    try:
        prs = Presentation(file_path)
    except Exception:
        return None

    sections: list[ReferenceSection] = []
    for i, slide in enumerate(prs.slides):
        title_text = ""
        body_texts: list[str] = []
        has_table = bool(slide.shapes.placeholders and
                         any(hasattr(sh, "table") for sh in slide.shapes))
        has_chart = any(sh.has_chart for sh in slide.shapes if hasattr(sh, "has_chart"))

        for shape in slide.shapes:
            if shape.has_text_frame:
                for ph in slide.placeholders:
                    if ph == shape and ph.placeholder_format.idx == 0:
                        title_text = shape.text_frame.text.strip()
                        break
                else:
                    text = shape.text_frame.text.strip()
                    if text and text != title_text:
                        body_texts.append(text)

        if not title_text:
            title_text = f"第{i + 1}页"

        char_count = sum(len(t) for t in body_texts)
        content_type = "table" if has_table else "mixed" if has_chart else "paragraphs"

        sections.append(ReferenceSection(
            heading_text=title_text,
            heading_level=1,
            content_type=content_type,
            char_count=char_count,
            has_table=has_table,
            has_bullets=False,
            paragraph_count=len(body_texts),
        ))

    total_chars = sum(s.char_count for s in sections)
    return ReferenceStructure(
        doc_type="slide_deck",
        sections=sections,
        style_note="structured_headings",
        estimated_total_chars=total_chars,
        slide_count=len(sections),
    )


# ── PDF extraction (best-effort) ─────────────────────────────────────────────

def extract_pdf_structure(file_path: str) -> Optional[ReferenceStructure]:
    """Extract heading structure from a PDF via heuristics on the extracted text."""
    try:
        from app.services.file_parser import extract_text
    except ImportError:
        return None

    try:
        text = extract_text(file_path)
    except Exception:
        return None

    if not text:
        return None

    lines = text.splitlines()
    sections: list[ReferenceSection] = []
    current_heading: Optional[str] = None
    current_lines: list[str] = []

    def _looks_like_heading(line: str) -> bool:
        line = line.strip()
        if not line or len(line) > 50:
            return False
        if _CN_HEADING_RE.match(line):
            return True
        if line.isupper() and len(line) > 3:
            return True
        return False

    def _flush():
        nonlocal current_heading, current_lines
        if current_heading:
            blob = " ".join(current_lines)
            sections.append(ReferenceSection(
                heading_text=current_heading,
                heading_level=1,
                content_type="paragraphs",
                char_count=len(blob),
                has_table=False,
                has_bullets=any(_BULLET_RE.match(l) for l in current_lines),
                paragraph_count=len(current_lines),
            ))
        current_heading = None
        current_lines = []

    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue
        if _looks_like_heading(stripped):
            _flush()
            current_heading = stripped
        elif current_heading:
            current_lines.append(stripped)

    _flush()

    if not sections:
        sections.append(ReferenceSection(
            heading_text="（全文）",
            heading_level=1,
            content_type="paragraphs",
            char_count=len(text),
            has_table=False,
            has_bullets=False,
            paragraph_count=len(lines),
        ))

    return ReferenceStructure(
        doc_type="report",
        sections=sections,
        style_note="structured_headings" if len(sections) > 1 else "narrative_prose",
        estimated_total_chars=sum(s.char_count for s in sections),
    )


# ── Dispatcher ───────────────────────────────────────────────────────────────

def extract_reference_structure(file_path: str) -> Optional[ReferenceStructure]:
    """Dispatch to the correct extractor based on file extension."""
    lower = file_path.lower()
    if lower.endswith(".docx") or lower.endswith(".doc"):
        return extract_docx_structure(file_path)
    elif lower.endswith(".pptx") or lower.endswith(".ppt"):
        return extract_pptx_structure(file_path)
    elif lower.endswith(".pdf"):
        return extract_pdf_structure(file_path)
    return None
