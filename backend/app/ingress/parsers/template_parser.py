"""模板占位符提取器 — 从 .dotx/.potx 提取占位符映射表.

Word (.dotx/.docx): 提取 {{变量名}} / {{#section}} / [PLACEHOLDER] 等模式
PowerPoint (.potx/.pptx): 提取每个 layout 的文本框占位符
"""
from __future__ import annotations

import io
import logging
import re
from dataclasses import dataclass, field
from typing import Any

logger = logging.getLogger(__name__)

# 常见占位符正则
_PLACEHOLDER_PATTERNS = [
    re.compile(r"\{\{([^}]+)\}\}"),          # {{变量名}}
    re.compile(r"\[([A-Z_][A-Z0-9_]{2,})\]"),  # [PLACEHOLDER]
    re.compile(r"<<([^>]+)>>"),               # <<字段名>>
    re.compile(r"\$\{([^}]+)\}"),             # ${变量名}
    re.compile(r"【([^】]{2,30})】"),           # 【字段名】
]


@dataclass
class PlaceholderInfo:
    name: str
    location: str        # "slide_N" / "section_header" / "body" etc.
    context_text: str    # 周边文本，帮助 LLM 理解用途
    data_type_hint: str  # 推断的数据类型: text/number/date/table/chart


@dataclass
class TemplateMeta:
    filename: str
    file_type: str              # docx/pptx/dotx/potx
    placeholders: list[PlaceholderInfo] = field(default_factory=list)
    slide_count: int = 0        # pptx only
    section_names: list[str] = field(default_factory=list)  # docx headings
    raw_text: str = ""

    def to_variable_map(self) -> dict[str, Any]:
        """返回 {占位符名: {location, hint}} 的填充映射表."""
        return {
            p.name: {"location": p.location, "type": p.data_type_hint, "context": p.context_text}
            for p in self.placeholders
        }

    def summary(self) -> str:
        lines = [f"模板文件: {self.filename} ({self.file_type})"]
        if self.slide_count:
            lines.append(f"幻灯片数: {self.slide_count}")
        if self.section_names:
            lines.append(f"章节: {', '.join(self.section_names[:10])}")
        lines.append(f"识别占位符 {len(self.placeholders)} 个:")
        for p in self.placeholders[:20]:
            lines.append(f"  [{p.location}] {{{{{p.name}}}}} — 类型:{p.data_type_hint}")
        return "\n".join(lines)


class TemplatePlaceholderParser:
    """从 Word/PPT 模板文件中提取占位符定义."""

    @classmethod
    def parse(cls, filename: str, content: bytes) -> TemplateMeta | None:
        lower = filename.lower()
        if lower.endswith((".docx", ".dotx", ".doc")):
            return cls._parse_docx(filename, content)
        elif lower.endswith((".pptx", ".potx", ".ppt")):
            return cls._parse_pptx(filename, content)
        return None

    @classmethod
    def _parse_docx(cls, filename: str, content: bytes) -> TemplateMeta:
        meta = TemplateMeta(filename=filename, file_type="docx")
        try:
            from docx import Document
            doc = Document(io.BytesIO(content))
            all_texts: list[tuple[str, str]] = []  # (text, location)

            for i, para in enumerate(doc.paragraphs):
                style = para.style.name if para.style else ""
                loc = "heading" if "Heading" in style else f"para_{i}"
                text = para.text.strip()
                if text:
                    all_texts.append((text, loc))
                if "Heading" in style and text:
                    meta.section_names.append(text)

            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        all_texts.append((cell.text.strip(), "table_cell"))

            full_text = "\n".join(t for t, _ in all_texts)
            meta.raw_text = full_text[:3000]

            seen: set[str] = set()
            for text, loc in all_texts:
                for pattern in _PLACEHOLDER_PATTERNS:
                    for m in pattern.finditer(text):
                        name = m.group(1).strip()
                        if name and name not in seen:
                            seen.add(name)
                            meta.placeholders.append(PlaceholderInfo(
                                name=name, location=loc,
                                context_text=text[:100],
                                data_type_hint=_infer_type(name),
                            ))
        except Exception as exc:
            logger.warning("TemplatePlaceholderParser docx failed: %s", exc)
        return meta

    @classmethod
    def _parse_pptx(cls, filename: str, content: bytes) -> TemplateMeta:
        meta = TemplateMeta(filename=filename, file_type="pptx")
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            meta.slide_count = len(prs.slides)
            seen: set[str] = set()

            for slide_idx, slide in enumerate(prs.slides):
                loc = f"slide_{slide_idx + 1}"
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for para in shape.text_frame.paragraphs:
                        text = "".join(r.text for r in para.runs).strip()
                        if not text:
                            continue
                        for pattern in _PLACEHOLDER_PATTERNS:
                            for m in pattern.finditer(text):
                                name = m.group(1).strip()
                                if name and name not in seen:
                                    seen.add(name)
                                    meta.placeholders.append(PlaceholderInfo(
                                        name=name, location=loc,
                                        context_text=text[:100],
                                        data_type_hint=_infer_type(name),
                                    ))

            # Also extract from slide layouts
            for layout in prs.slide_layouts:
                for ph in layout.placeholders:
                    name = ph.name
                    if name and name not in seen:
                        seen.add(name)
                        meta.placeholders.append(PlaceholderInfo(
                            name=name, location="layout",
                            context_text=f"Layout placeholder: {ph.placeholder_format.type}",
                            data_type_hint="text",
                        ))
        except Exception as exc:
            logger.warning("TemplatePlaceholderParser pptx failed: %s", exc)
        return meta


def _infer_type(name: str) -> str:
    """从占位符名称推断数据类型."""
    low = name.lower()
    if any(k in low for k in ("date", "日期", "time", "年", "月", "日")):
        return "date"
    if any(k in low for k in ("rate", "ratio", "pct", "percent", "率", "比", "%", "增长")):
        return "number"
    if any(k in low for k in ("amount", "total", "sum", "金额", "合计", "数量", "count")):
        return "number"
    if any(k in low for k in ("chart", "graph", "图表", "图", "趋势")):
        return "chart"
    if any(k in low for k in ("table", "表格", "列表")):
        return "table"
    return "text"
