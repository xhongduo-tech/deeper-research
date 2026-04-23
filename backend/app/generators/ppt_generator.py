import io
import os
from typing import Optional, List, Dict, Any

# Color constants
PRIMARY_COLOR = (0x1F, 0x49, 0x7D)   # Deep blue
ACCENT_COLOR = (0x2E, 0x75, 0xB6)    # Medium blue
LIGHT_BG = (0xF0, 0xF4, 0xF8)        # Light blue-gray
WHITE = (0xFF, 0xFF, 0xFF)
DARK_TEXT = (0x1A, 0x1A, 0x2E)
GRAY_TEXT = (0x44, 0x44, 0x44)


class PPTGenerator:
    """
    Generates professional PowerPoint presentations.
    Supports custom templates and creates clean designs from scratch.
    """

    async def generate(
        self,
        content: dict,
        template_path: Optional[str] = None,
        theme: str = "professional_blue",
    ) -> bytes:
        """
        Generate a PPTX file.

        Args:
            content: {
                "title": "Presentation Title",
                "subtitle": "Optional subtitle",
                "author": "Author name",
                "date": "Date string",
                "sections": [
                    {
                        "title": "Section Title",
                        "content": "Main text content",
                        "bullet_points": ["Point 1", "Point 2"],
                        "charts": [],  # matplotlib figure bytes list
                        "table": [[row1], [row2]],  # optional table data
                        "notes": "Speaker notes",
                        "layout": "title_content" | "two_column" | "blank"
                    }
                ],
                "summary": "Optional summary slide content"
            }
            template_path: Path to a .pptx template file
            theme: Color theme name

        Returns:
            bytes of the .pptx file
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN

            if template_path and os.path.exists(template_path):
                prs = Presentation(template_path)
            else:
                prs = Presentation()

            # Set slide dimensions (widescreen 16:9)
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)

            # Title slide
            self._add_title_slide(prs, content)

            # Content slides
            for section in content.get("sections", []):
                layout = section.get("layout", "title_content")
                if layout == "two_column":
                    self._add_two_column_slide(prs, section)
                elif section.get("table"):
                    self._add_table_slide(prs, section)
                elif section.get("bullet_points"):
                    self._add_bullet_slide(prs, section)
                else:
                    self._add_content_slide(prs, section)

            # Summary slide if provided
            if content.get("summary"):
                self._add_summary_slide(prs, content["summary"], content.get("title", ""))

            # Save to bytes
            output = io.BytesIO()
            prs.save(output)
            output.seek(0)
            return output.read()

        except ImportError:
            return self._generate_text_fallback(content)
        except Exception as e:
            raise RuntimeError(f"PPT生成失败: {str(e)}")

    def _add_title_slide(self, prs, content: dict):
        """Add the title/cover slide."""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*PRIMARY_COLOR)

        # Title text box
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(11.33), Inches(2)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = 1  # CENTER
        run = p.add_run()
        run.text = content.get("title", "报告标题")
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.color.rgb = RGBColor(*WHITE)

        # Subtitle
        if content.get("subtitle"):
            sub_box = slide.shapes.add_textbox(
                Inches(1), Inches(4.2), Inches(11.33), Inches(0.8)
            )
            tf2 = sub_box.text_frame
            p2 = tf2.paragraphs[0]
            p2.alignment = 1
            run2 = p2.add_run()
            run2.text = content["subtitle"]
            run2.font.size = Pt(20)
            run2.font.color.rgb = RGBColor(0xB0, 0xC8, 0xE8)

        # Author and date
        info_parts = []
        if content.get("author"):
            info_parts.append(content["author"])
        if content.get("date"):
            info_parts.append(content["date"])

        if info_parts:
            info_box = slide.shapes.add_textbox(
                Inches(1), Inches(5.5), Inches(11.33), Inches(0.6)
            )
            tf3 = info_box.text_frame
            p3 = tf3.paragraphs[0]
            p3.alignment = 1
            run3 = p3.add_run()
            run3.text = "  |  ".join(info_parts)
            run3.font.size = Pt(14)
            run3.font.color.rgb = RGBColor(0xB0, 0xC8, 0xE8)

        # Bottom accent line
        line = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(0), Inches(6.9), Inches(13.33), Inches(0.6)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(*ACCENT_COLOR)
        line.line.fill.background()

    def _add_content_slide(self, prs, section: dict):
        """Add a standard content slide."""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # White background
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(*WHITE)

        # Header bar
        header = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(13.33), Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(*PRIMARY_COLOR)
        header.line.fill.background()

        # Title in header
        title_box = slide.shapes.add_textbox(
            Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.9)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = section.get("title", "")
        run.font.size = Pt(26)
        run.font.bold = True
        run.font.color.rgb = RGBColor(*WHITE)

        # Content area
        content = section.get("content", "")
        if content:
            content_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.5)
            )
            tf2 = content_box.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            run2 = p2.add_run()
            run2.text = content[:1500]  # Limit content
            run2.font.size = Pt(14)
            run2.font.color.rgb = RGBColor(*DARK_TEXT)

        # Speaker notes
        if section.get("notes"):
            slide.notes_slide.notes_text_frame.text = section["notes"]

    def _add_bullet_slide(self, prs, section: dict):
        """Add a slide with bullet points."""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(*WHITE)

        # Header
        header = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(13.33), Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(*PRIMARY_COLOR)
        header.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.9)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = section.get("title", "")
        run.font.size = Pt(26)
        run.font.bold = True
        run.font.color.rgb = RGBColor(*WHITE)

        # Intro text if any
        y_start = 1.4
        if section.get("content"):
            intro_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(y_start), Inches(12.3), Inches(0.8)
            )
            tf_i = intro_box.text_frame
            tf_i.word_wrap = True
            p_i = tf_i.paragraphs[0]
            run_i = p_i.add_run()
            run_i.text = section["content"][:300]
            run_i.font.size = Pt(13)
            run_i.font.color.rgb = RGBColor(*GRAY_TEXT)
            y_start = 2.3

        # Bullet points
        bullets = section.get("bullet_points", [])
        if bullets:
            bullet_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(y_start), Inches(12.3), Inches(7.5 - y_start - 0.3)
            )
            tf_b = bullet_box.text_frame
            tf_b.word_wrap = True

            for i, bullet in enumerate(bullets[:10]):  # Max 10 bullets
                if i == 0:
                    p = tf_b.paragraphs[0]
                else:
                    p = tf_b.add_paragraph()

                p.level = 0
                run = p.add_run()
                run.text = f"• {bullet}"
                run.font.size = Pt(15)
                run.font.color.rgb = RGBColor(*DARK_TEXT)

                # Add spacing
                from pptx.util import Pt as PtU
                p.space_after = PtU(6)

        if section.get("notes"):
            slide.notes_slide.notes_text_frame.text = section["notes"]

    def _add_table_slide(self, prs, section: dict):
        """Add a slide with a data table."""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(*WHITE)

        # Header
        header = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(13.33), Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(*PRIMARY_COLOR)
        header.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.9)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = section.get("title", "")
        run.font.size = Pt(26)
        run.font.bold = True
        run.font.color.rgb = RGBColor(*WHITE)

        # Table
        table_data = section.get("table", [])
        if table_data and len(table_data) > 0:
            rows = min(len(table_data), 15)
            cols = min(len(table_data[0]), 8) if table_data else 3

            table = slide.shapes.add_table(
                rows, cols,
                Inches(0.3), Inches(1.4),
                Inches(12.7), Inches(5.8)
            ).table

            for row_idx in range(rows):
                for col_idx in range(cols):
                    cell = table.cell(row_idx, col_idx)
                    if row_idx < len(table_data) and col_idx < len(table_data[row_idx]):
                        cell.text = str(table_data[row_idx][col_idx])
                    else:
                        cell.text = ""

                    # Style header row
                    tf_c = cell.text_frame
                    p_c = tf_c.paragraphs[0]
                    run_c = p_c.runs[0] if p_c.runs else p_c.add_run()
                    run_c.font.size = Pt(11)
                    if row_idx == 0:
                        run_c.font.bold = True
                        run_c.font.color.rgb = RGBColor(*WHITE)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(*ACCENT_COLOR)
                    elif row_idx % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(*LIGHT_BG)

    def _add_two_column_slide(self, prs, section: dict):
        """Add a two-column layout slide."""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(*WHITE)

        # Header
        header = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(13.33), Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(*PRIMARY_COLOR)
        header.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.9)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = section.get("title", "")
        run.font.size = Pt(26)
        run.font.bold = True
        run.font.color.rgb = RGBColor(*WHITE)

        # Left column
        left_content = section.get("left_content", section.get("content", ""))
        left_box = slide.shapes.add_textbox(
            Inches(0.4), Inches(1.4), Inches(6.0), Inches(5.7)
        )
        tf_l = left_box.text_frame
        tf_l.word_wrap = True
        p_l = tf_l.paragraphs[0]
        run_l = p_l.add_run()
        run_l.text = left_content[:1000]
        run_l.font.size = Pt(13)
        run_l.font.color.rgb = RGBColor(*DARK_TEXT)

        # Divider
        divider = slide.shapes.add_shape(
            1, Inches(6.6), Inches(1.4), Inches(0.05), Inches(5.5)
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = RGBColor(*ACCENT_COLOR)
        divider.line.fill.background()

        # Right column
        right_content = section.get("right_content", "")
        if right_content:
            right_box = slide.shapes.add_textbox(
                Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.7)
            )
            tf_r = right_box.text_frame
            tf_r.word_wrap = True
            p_r = tf_r.paragraphs[0]
            run_r = p_r.add_run()
            run_r.text = right_content[:1000]
            run_r.font.size = Pt(13)
            run_r.font.color.rgb = RGBColor(*DARK_TEXT)

    def _add_summary_slide(self, prs, summary: str, title: str):
        """Add a conclusion/summary slide."""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Dark gradient background
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(*PRIMARY_COLOR)

        # "CONCLUSION" label
        label_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.5)
        )
        tf_label = label_box.text_frame
        p_label = tf_label.paragraphs[0]
        run_label = p_label.add_run()
        run_label.text = "结论与展望"
        run_label.font.size = Pt(14)
        run_label.font.color.rgb = RGBColor(0xB0, 0xC8, 0xE8)
        run_label.font.bold = True

        # Summary title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.2)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"{title} - 综合总结"
        run.font.size = Pt(30)
        run.font.bold = True
        run.font.color.rgb = RGBColor(*WHITE)

        # Summary content
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.8), Inches(12.3), Inches(4.0)
        )
        tf2 = content_box.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = summary[:1200]
        run2.font.size = Pt(16)
        run2.font.color.rgb = RGBColor(0xD0, 0xE8, 0xFF)

        # Bottom accent
        accent = slide.shapes.add_shape(
            1, Inches(0.5), Inches(6.8), Inches(4.0), Inches(0.08)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(0x70, 0xB8, 0xFF)
        accent.line.fill.background()

    def _generate_text_fallback(self, content: dict) -> bytes:
        """Fallback when python-pptx is not available."""
        text = f"# {content.get('title', 'Report')}\n\n"
        for section in content.get("sections", []):
            text += f"## {section.get('title', '')}\n\n"
            text += section.get("content", "") + "\n\n"
            for bp in section.get("bullet_points", []):
                text += f"- {bp}\n"
        return text.encode("utf-8")

    async def generate_from_markdown(self, markdown: str, title: str = "Report") -> bytes:
        """
        Parse a markdown document and convert it to PPTX.
        Each H2 header becomes a new slide.
        """
        import re
        sections = []
        current_section = None
        current_bullets = []
        current_content = []

        for line in markdown.split("\n"):
            if line.startswith("## "):
                if current_section:
                    sections.append({
                        "title": current_section,
                        "content": "\n".join(current_content).strip(),
                        "bullet_points": current_bullets,
                    })
                current_section = line[3:].strip()
                current_content = []
                current_bullets = []
            elif line.startswith("- ") or line.startswith("* "):
                current_bullets.append(line[2:].strip())
            elif line.startswith("# "):
                title = line[2:].strip()
            elif current_section and line.strip():
                current_content.append(line)

        if current_section:
            sections.append({
                "title": current_section,
                "content": "\n".join(current_content).strip(),
                "bullet_points": current_bullets,
            })

        content = {
            "title": title,
            "sections": sections,
        }
        return await self.generate(content)
