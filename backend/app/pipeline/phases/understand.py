"""UNDERSTAND phase — parses user intent and extracts reference document structure.

Key fixes vs. the old _phase_understand:
1. Intent is simplified to 3 values: fill_from_reference | fresh | extend
   (removes the ambiguous "format_only" / "mixed" values)
2. intent_confidence < 0.7 → force clarification (no longer optional)
3. Reference doc structure is extracted DETERMINISTICALLY (no LLM) using
   ref_extractor and stored in scoping_plan for use by the PLAN phase.
"""
from __future__ import annotations

import logging
from typing import TYPE_CHECKING

from app.pipeline.types import PipelineContext, PipelineError
from app.pipeline.llm_helpers import call_llm_json, _coerce_bool, _listify
from app.pipeline.skills_loader import (
    build_skill_context, extract_skills_from_brief, filter_skills_for_phase
)
from app.pipeline.claim_utils import extract_numeric_baseline, parse_branch_context
from app.rendering.ref_extractor import extract_reference_structure

if TYPE_CHECKING:
    pass

logger = logging.getLogger(__name__)

_UNDERSTAND_SYSTEM = """你是需求解析专家。分析用户的文档生成需求，提取核心意图。
只输出JSON，不附加任何解释文字。

【intent 识别规则】
- "fill_from_reference"：用户希望基于上传的参考文档直接生成当前版本
  触发条件：briefs 中出现"根据X生成Y""仿照参考文档""续写""下一年度同类文档"等
- "extend"：用户在参考文档基础上延伸，保留部分内容并增加新内容
  触发条件："补充""扩充""增加章节""在此基础上"等
- "fresh"：全新内容，不受参考文档约束（无上传文件，或明确说不参考格式）
  触发条件：无上传文件，或明确说"全新生成"

【tense 识别规则】
- 述职报告/工作总结/年终汇报 = "past"（已完成的工作，使用过去时）
- 计划/方案/预测 = "present"

【style 识别规则】
- 演讲稿/述职稿（连续散文、无正式标题） = "speech"
- 正式报告/白皮书/研究报告 = "report"
- 数据分析/经营分析 = "analysis"
- 表格密集文档 = "table_heavy"

仅当关键信息缺失时（时间范围、数据口径、明确对比主体）将clarification_needed设为true。"""


class UnderstandPhase:
    PHASE_NAME = "UNDERSTAND"

    def __init__(self, ctx: PipelineContext):
        self.ctx = ctx

    async def run(self) -> None:
        ctx = self.ctx

        # ── P2-F: Resume from cache when user has answered clarifications ─────
        scoping = ctx.report.scoping_plan or {}
        if scoping.get("understanding") and scoping.get("answered_clarifications"):
            understanding = dict(scoping["understanding"])
            answers_text = "; ".join(
                f"{a.get('question', '')}: {a.get('answer', '')}"
                for a in scoping["answered_clarifications"]
            )
            understanding["clarification_needed"] = False
            understanding["clarification_questions"] = []
            understanding["clarification_answers"] = answers_text
            ctx.understanding = understanding
            logger.info(
                "[UNDERSTAND] Restored from cache with %d clarification answers",
                len(scoping["answered_clarifications"]),
            )
            return

        brief = ctx.brief
        report_type = ctx.report_type
        output_format = ctx.output_format
        uploaded_texts = ctx.uploaded_texts

        # ── LLM-OS: Intent Router + Ingress asset enrichment ─────────────────
        await self._run_ingress_and_routing(ctx, brief, output_format, uploaded_texts)

        # ── Deterministic pre-processing (no LLM) ────────────────────────────
        numeric_baseline = extract_numeric_baseline(uploaded_texts)
        branch_context = parse_branch_context(uploaded_texts)
        reference_structure = await self._extract_reference_structure()

        # Auto-detect skills
        skill_names = list(ctx.skills or extract_skills_from_brief(brief))
        if reference_structure and "reference-style-miner" not in skill_names:
            skill_names.append("reference-style-miner")

        skill_context = build_skill_context(filter_skills_for_phase(skill_names, "understand"))

        # ── LLM: intent classification ────────────────────────────────────────
        format_desc = _format_description(output_format)
        evidence_summary = _build_evidence_summary(uploaded_texts)
        baseline_summary = _build_baseline_summary(numeric_baseline)

        system_prompt = _UNDERSTAND_SYSTEM
        if skill_context:
            system_prompt += f"\n\n以下是你必须遵循的专业技能规范（Skill）：\n\n{skill_context}"

        ref_hint = ""
        if reference_structure:
            ref_hint = (
                f"\n\n【参考文档结构已检测】文档类型：{reference_structure.doc_type}，"
                f"共 {len(reference_structure.sections)} 个章节，"
                f"风格：{reference_structure.style_note}。"
                f"前5章：{', '.join(s.heading_text[:20] for s in reference_structure.sections[:5])}"
            )

        user_prompt = f"""请分析以下文档生成需求：

【需求描述】
{brief}

【报告类型】{report_type}
【输出格式】{format_desc}{ref_hint}

【上传材料摘要】
{evidence_summary or "（无上传材料）"}

【上传材料中的数值基线】
{baseline_summary or "（未提取到数值基线）"}

请输出严格JSON：
{{
  "topic": "主题（1句话概括）",
  "audience": "目标受众",
  "goal": "用户核心目标",
  "format_hint": "格式细节偏好",
  "key_questions": ["问题1", "问题2"],
  "tone": "专业|简洁|深入|通俗",
  "estimated_sections": 6,
  "intent": "fill_from_reference|extend|fresh",
  "intent_confidence": 0.9,
  "tense": "past|present",
  "style": "speech|report|analysis|table_heavy",
  "clarification_needed": false,
  "clarification_questions": []
}}"""

        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ]

        result = await call_llm_json(
            messages,
            temperature=0.3,
            tier="heavy",
            fallback={
                "topic": brief[:50],
                "audience": "业务决策者",
                "goal": brief[:100],
                "format_hint": format_desc,
                "key_questions": [],
                "tone": "专业",
                "estimated_sections": 6,
                "intent": "fresh",
                "intent_confidence": 0.5,
                "tense": "present",
                "style": "report",
                "clarification_needed": False,
                "clarification_questions": [],
            },
        )

        # P1-1: Infer narrative framework (SCR / Problem-Solution / Data-Driven / Chronological)
        narrative_framework = _infer_narrative_framework(
            brief=brief,
            style=result.get("style", "report"),
            tense=result.get("tense", "present"),
            topic=str(result.get("topic", "")),
        )

        # P1-4: Build entity registry (canonical org names, metric names, units)
        entity_registry = _build_entity_registry(uploaded_texts, brief)

        # ── Normalize and validate intent ─────────────────────────────────────
        intent = str(result.get("intent", "fresh")).strip()
        if intent not in ("fill_from_reference", "extend", "fresh"):
            # Map old values
            _intent_map = {
                "fill_from_reference": "fill_from_reference",
                "format_only": "fresh",
                "mixed": "extend",
            }
            intent = _intent_map.get(intent, "fresh")

        intent_confidence = float(result.get("intent_confidence", 0.5) or 0.5)
        tense = result.get("tense", "present")
        if tense not in ("past", "present"):
            tense = "present"
        style = result.get("style", "report")
        if style not in ("speech", "report", "analysis", "table_heavy"):
            style = "report"

        clarification_needed = _coerce_bool(result.get("clarification_needed"), False)
        # Force clarification on low-confidence intent
        if intent_confidence < 0.7 and not clarification_needed and not getattr(ctx, "skip_clarify", False):
            clarification_needed = True
            result.setdefault("clarification_questions", [])
            result["clarification_questions"].append({
                "question": f"您的意图是否是基于已上传的参考文档生成当前版本？（当前识别：{intent}，置信度：{intent_confidence:.0%}）",
                "default_answer": "是，请基于参考文档生成",
                "priority": "high",
            })
        if getattr(ctx, "skip_clarify", False):
            clarification_needed = False
            result["clarification_questions"] = []

        understanding = {
            "topic": str(result.get("topic", brief[:50])),
            "audience": str(result.get("audience", "业务决策者")),
            "goal": str(result.get("goal", brief[:100])),
            "format_hint": str(result.get("format_hint", format_desc)),
            "key_questions": [str(q) for q in _listify(result.get("key_questions", []))],
            "tone": str(result.get("tone", "专业")),
            "estimated_sections": int(result.get("estimated_sections", 6) or 6),
            "intent": intent,
            "intent_confidence": intent_confidence,
            "tense": tense,
            "style": style,
            "clarification_needed": clarification_needed,
            "clarification_questions": _listify(result.get("clarification_questions", [])),
            "numeric_baseline": numeric_baseline,
            "active_skills": skill_names,
            "narrative_framework": narrative_framework,   # P1-1
            "entity_registry": entity_registry,           # P1-4
        }
        if branch_context:
            understanding["branch_context"] = branch_context
        if reference_structure:
            understanding["reference_structure"] = reference_structure.to_dict()
        # P2-4: Store extracted template DNA (set by _extract_reference_structure)
        template_dna = getattr(self, "_template_dna", None)
        if template_dna:
            understanding["template_dna"] = template_dna

        ctx.understanding = understanding

        # Persist to DB
        ctx.report.scoping_plan = ctx.report.scoping_plan or {}
        ctx.report.scoping_plan["understanding"] = understanding
        await ctx.db.commit()

        # Handle clarification
        if clarification_needed and result.get("clarification_questions"):
            await self._create_clarifications(result["clarification_questions"])

    async def _run_ingress_and_routing(
        self,
        ctx,
        brief: str,
        output_format: str,
        uploaded_texts: list[str],
    ) -> None:
        """LLM-OS: 统一接入解析 + 意图路由（非阻塞，失败不影响主流程）."""
        import uuid
        try:
            from sqlalchemy import select
            from app.models.uploaded_file import UploadedFile
            from app.ingress.vfs import VirtualFileSystem
            from app.ingress.dispatcher import IngressDispatcher
            from app.knowledge.intent_router import IntentRouter

            # 1) 取所有已上传文件
            rows = (await ctx.db.execute(
                select(UploadedFile).where(UploadedFile.report_id == ctx.report.id)
            )).scalars().all()

            file_types: list[str] = []
            all_assets = []
            vfs_trees: list[str] = []

            for row in rows:
                from pathlib import Path, PurePosixPath
                suffix = PurePosixPath(row.original_name).suffix.lower()
                file_types.append(suffix)

                archive_exts = {".zip", ".tar.gz", ".tgz", ".tar.bz2", ".tar"}
                is_archive = any(row.original_name.lower().endswith(ext) for ext in archive_exts)

                try:
                    content = Path(row.file_path).read_bytes()
                    vfs = VirtualFileSystem.from_bytes(row.original_name, content)
                    if is_archive:
                        vfs_trees.append(vfs.directory_tree())
                    assets = await IngressDispatcher.dispatch_vfs(vfs)
                    all_assets.extend(assets)

                    # 结构化数据文件 → 注册到 DuckDB
                    if suffix in (".xlsx", ".xls", ".csv", ".xlsb"):
                        if not ctx.duckdb_session_id:
                            ctx.duckdb_session_id = f"report_{ctx.report.id}_{uuid.uuid4().hex[:8]}"
                        from app.compute.duckdb_engine import DuckDBEngine
                        engine = DuckDBEngine.get_or_create(ctx.duckdb_session_id)
                        try:
                            engine.register_file(row.original_name, content)
                            ctx.duckdb_schema = engine.table_schema_text()
                        except Exception as exc:
                            logger.warning("[UNDERSTAND] DuckDB register failed: %s", exc)

                    # 模板文件 → 提取占位符
                    if row.is_template and suffix in (".dotx", ".potx", ".docx", ".pptx"):
                        from app.ingress.parsers.template_parser import TemplatePlaceholderParser
                        meta = TemplatePlaceholderParser.parse(row.original_name, content)
                        if meta:
                            ctx.template_meta = meta

                    # 补充 uploaded_texts（代码/config 的 context_text）
                    for asset in assets:
                        if asset.context_text and asset.asset_type in ("code", "config"):
                            uploaded_texts.append(asset.context_text)

                except Exception as exc:
                    logger.warning("[UNDERSTAND] Ingress failed for %s: %s", row.original_name, exc)

            ctx.ingress_assets = all_assets
            if vfs_trees:
                ctx.vfs_tree = "\n\n".join(vfs_trees)

            # 2) 意图路由
            intent = await IntentRouter.route(
                brief=brief,
                file_types=file_types,
                output_format=output_format,
            )
            ctx.intent = intent
            logger.info("[UNDERSTAND] Intent: %s (conf=%.2f)", intent.scenario, intent.confidence)

        except Exception as exc:
            logger.warning("[UNDERSTAND] Ingress/routing step failed (non-fatal): %s", exc)

    async def _extract_reference_structure(self):
        """Deterministically extract structure from the template file (if any).

        P2-4: Also extracts template DNA (accent color + font) from PPTX templates
        and stores it in ctx.understanding["template_dna"] for SPEC_GEN injection.
        """
        from sqlalchemy import select
        from app.models.uploaded_file import UploadedFile

        result = await self.ctx.db.execute(
            select(UploadedFile).where(
                UploadedFile.report_id == self.ctx.report.id,
                UploadedFile.is_template == True,
            )
        )
        template_file = result.scalar_one_or_none()

        if template_file and template_file.file_path:
            try:
                structure = extract_reference_structure(template_file.file_path)
                if structure:
                    logger.info(
                        "[UNDERSTAND] ref structure: %d sections, type=%s",
                        len(structure.sections), structure.doc_type,
                    )

                # P2-4: Extract template DNA for PPTX files (non-fatal)
                fpath = template_file.file_path
                if fpath.lower().endswith((".pptx", ".ppt")):
                    try:
                        dna = _extract_pptx_dna(fpath)
                        if dna:
                            # Merge into understanding (will be set later in run())
                            self._template_dna = dna
                            logger.info("[UNDERSTAND] PPTX template DNA: %s", dna)
                    except Exception as dna_exc:
                        logger.debug("[UNDERSTAND] PPTX DNA extraction failed: %s", dna_exc)

                return structure
            except Exception as exc:
                logger.warning("[UNDERSTAND] ref structure extraction failed: %s", exc)
        return None

    async def _create_clarifications(self, questions: list) -> None:
        from app.models.clarification import Clarification
        for q_item in questions[:3]:
            if isinstance(q_item, dict):
                question_text = q_item.get("question", "")
                default = q_item.get("default_answer", "")
                priority = q_item.get("priority", "medium")
            else:
                question_text = str(q_item)
                default = ""
                priority = "medium"
            if not question_text:
                continue
            clarification = Clarification(
                report_id=self.ctx.report.id,
                question=question_text,
                default_answer=default,
                priority=priority,
                status="pending",
            )
            self.ctx.db.add(clarification)
        await self.ctx.db.commit()


def _extract_pptx_dna(pptx_path: str) -> dict:
    """P2-4: Extract accent color and font name from a PPTX template file.

    Returns dict with keys 'accent_hex' (6-char hex, no #) and 'font_name'.
    Used by SPEC_GEN to adapt slide content style to the template's visual design.
    """
    try:
        from pptx import Presentation
        from pptx.util import Pt
        prs = Presentation(pptx_path)

        # Extract theme accent color (accent1 is the primary accent)
        accent_hex = ""
        try:
            from pptx.dml.color import RGBColor
            theme = prs.slide_master.theme_color_map
            # Try theme element directly via XML
            from lxml import etree
            ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
            theme_el = prs.slide_master._element.find(
                f".//{{{ns}}}theme/{{{ns}}}themeElements/{{{ns}}}clrScheme/{{{ns}}}accent1"
            )
            if theme_el is not None:
                srgb = theme_el.find(f"{{{ns}}}srgbClr")
                if srgb is not None:
                    accent_hex = srgb.get("val", "")
        except Exception:
            pass

        # Simpler fallback: look at first slide's title placeholder color
        if not accent_hex and prs.slides:
            try:
                for shape in prs.slides[0].shapes:
                    if shape.has_text_frame and shape.text_frame.paragraphs:
                        for run in shape.text_frame.paragraphs[0].runs:
                            if run.font.color and run.font.color.type:
                                c = run.font.color.rgb
                                if c:
                                    accent_hex = str(c)
                                    break
                        if accent_hex:
                            break
            except Exception:
                pass

        # Extract font name from slide master body font
        font_name = ""
        try:
            body_font = prs.slide_master.slide_layout_master.body_text_frame
        except Exception:
            body_font = None
        if not font_name:
            try:
                theme_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
                font_scheme = prs.slide_master._element.find(
                    f".//{{{theme_ns}}}theme/{{{theme_ns}}}themeElements/{{{theme_ns}}}fontScheme/{{{theme_ns}}}majorFont/{{{theme_ns}}}latin"
                )
                if font_scheme is not None:
                    font_name = font_scheme.get("typeface", "")
            except Exception:
                pass

        return {k: v for k, v in {"accent_hex": accent_hex, "font_name": font_name}.items() if v}
    except Exception:
        return {}


def _format_description(output_format: str) -> str:
    return {
        "ppt": "PPT演示文稿，断言式标题、要点精炼、数字优先",
        "pptx": "PPT演示文稿，断言式标题、要点精炼、数字优先",
        "word": "Word深度报告，结构化论述、数据支撑、逻辑递进",
        "doc": "Word深度报告，结构化论述、数据支撑、逻辑递进",
        "docx": "Word深度报告，结构化论述、数据支撑、逻辑递进",
        "excel": "Excel数据分析，指标提取、数据表格、量化计算",
        "xlsx": "Excel数据分析，指标提取、数据表格、量化计算",
    }.get(output_format.lower(), "文档")


def _build_evidence_summary(uploaded_texts: list[str] | None) -> str:
    if not uploaded_texts:
        return ""
    excerpts = [src[:3000] for src in uploaded_texts[:8]]
    return "\n\n---\n\n".join(excerpts)[:10000]


def _build_baseline_summary(numeric_baseline: dict) -> str:
    if not numeric_baseline:
        return ""
    items = []
    for key, info in list(numeric_baseline.items())[:20]:
        nums = ", ".join(info.get("numbers", [])[:4])
        items.append(f"- {key[:60]}（数值：{nums}，来源：{info.get('source', '')}）")
    return "\n".join(items)


def _infer_narrative_framework(brief: str, style: str, tense: str, topic: str) -> str:
    """P1-1: Rule-based narrative framework inference for PLAN phase injection.

    Returns one of: SCR | Problem-Solution | Data-Driven | Chronological
    """
    text = f"{brief} {topic}".lower()
    problem_kws = {"问题", "挑战", "根因", "原因", "解决", "方案", "改进", "优化", "整改"}
    data_kws = {"数据", "分析", "指标", "趋势", "同比", "环比", "增速", "占比", "占", "排名"}
    chrono_kws = {"总结", "述职", "回顾", "历程", "时间线", "进度", "阶段", "过程"}

    if any(kw in text for kw in problem_kws):
        return "Problem-Solution"
    if style == "analysis" or any(kw in text for kw in data_kws):
        return "Data-Driven"
    if tense == "past" or any(kw in text for kw in chrono_kws):
        return "Chronological"
    return "SCR"


def _build_entity_registry(uploaded_texts: list[str] | None, brief: str) -> dict:
    """P1-4: Extract canonical entity names and units from uploaded texts.

    Returns {entity_name: {"type": "org|metric", "unit": str}} for injection
    into all SPEC_GEN prompts so names/units stay consistent across sections.
    """
    import re
    registry: dict = {}
    sample = brief + "\n" + "\n".join(t[:2000] for t in (uploaded_texts or [])[:5])

    # Organisation names (e.g. XX公司, XX集团, XX部门)
    for m in re.finditer(r"[^\s，。、]{2,15}(?:公司|集团|部门|单位|中心|事业部|子公司)", sample):
        registry[m.group()] = {"type": "org", "unit": ""}

    # Metric names with optional unit (e.g. 营收额 100亿元, 增长率 15%)
    for m in re.finditer(
        r"([一-龥a-zA-Z]{2,10}(?:率|额|量|数|收入|利润|成本|费用|增速))"
        r"[为：:是]?\s*[\d,\.]+\s*(%|万|亿|元|人|个|次|倍)?",
        sample,
    ):
        registry[m.group(1)] = {"type": "metric", "unit": m.group(2) or ""}

    return dict(list(registry.items())[:20])
