"""QA phase — validates the rendered document against spec and numeric baseline.

Validates the ACTUAL rendered file (not intermediate markdown) to catch real
output quality issues. P0 issues trigger targeted SPEC_GEN re-run for failing
sections, then a second DOC_RENDER pass.

After one revision cycle, remaining failures are raised as PipelineError
rather than being silently degraded.
"""
from __future__ import annotations

import logging
import re
from typing import TYPE_CHECKING

from app.pipeline.types import PipelineContext, PipelineError
from app.pipeline.claim_utils import verify_section_claims

if TYPE_CHECKING:
    pass

logger = logging.getLogger(__name__)

_MIN_TOTAL_CHARS_ABSOLUTE = 200    # hard floor regardless of spec
_MIN_TOTAL_CHARS_RATIO = 0.30      # at least 30% of sum(target_chars) must be present
_SLIDE_BODY_CHAR_CAPACITY = 300    # approximate chars that fit in a slide body at 18pt


class QAPhase:
    PHASE_NAME = "QA"

    def __init__(self, ctx: PipelineContext):
        self.ctx = ctx

    async def run(self) -> None:
        ctx = self.ctx
        if not ctx.rendered_bytes:
            raise PipelineError(self.PHASE_NAME, "No rendered bytes to validate")

        issues = await self._check(ctx.rendered_bytes)
        p0_issues = [i for i in issues if i.get("severity") == "p0"]

        if not p0_issues:
            logger.info("[QA] Passed (%d total issues, 0 P0)", len(issues))
            return

        logger.warning("[QA] %d P0 issues found — triggering targeted revision", len(p0_issues))

        # Group by section_id
        failed_sections = list({i.get("section_id") for i in p0_issues if i.get("section_id")})

        if failed_sections and ctx.spec is not None:
            # Re-generate only the failing sections
            from app.pipeline.phases.spec_gen import SpecGenPhase
            from app.pipeline.phases.doc_render import DocRenderPhase

            try:
                spec_gen = SpecGenPhase(ctx, qa_feedback=p0_issues)
                await spec_gen.run(section_ids=failed_sections)
                doc_render = DocRenderPhase(ctx)
                await doc_render.run()
            except PipelineError:
                raise
            except Exception as exc:
                raise PipelineError(self.PHASE_NAME, f"Revision pass failed: {exc}") from exc

            # Second check
            issues2 = await self._check(ctx.rendered_bytes)
            p0_issues2 = [i for i in issues2 if i.get("severity") == "p0"]
            if p0_issues2:
                summary = "; ".join(i.get("message", "") for i in p0_issues2[:3])
                raise PipelineError(
                    self.PHASE_NAME,
                    f"QA failed after revision. Remaining P0 issues: {summary}",
                )
        else:
            # No section-level failures — just report-level
            summary = "; ".join(i.get("message", "") for i in p0_issues[:3])
            raise PipelineError(self.PHASE_NAME, f"QA failed: {summary}")

        logger.info("[QA] Passed after revision")

    async def _check(self, rendered_bytes: bytes) -> list[dict]:
        """Run deterministic checks on the rendered document bytes."""
        issues: list[dict] = []

        # Extract text from rendered bytes
        text = _extract_text_from_bytes(rendered_bytes, self.ctx.rendered_ext or "docx")
        clean_text = re.sub(r"\s+", "", text)

        # Check 1: total character count — dynamic floor based on spec target
        min_chars = _MIN_TOTAL_CHARS_ABSOLUTE
        if self.ctx.spec:
            target_total = _sum_target_chars(self.ctx.spec)
            if target_total > 0:
                min_chars = max(_MIN_TOTAL_CHARS_ABSOLUTE,
                                int(target_total * _MIN_TOTAL_CHARS_RATIO))
        if len(clean_text) < min_chars:
            issues.append({
                "severity": "p0",
                "check": "total_chars",
                "message": (
                    f"文档内容不足（实际 {len(clean_text)} 字，"
                    f"要求至少 {min_chars} 字），可能生成失败"
                ),
                "section_id": None,
            })

        if not self.ctx.spec:
            return issues

        # Check 2: section presence
        for sec in _iter_sections(self.ctx.spec):
            sec_id = sec.get("id", "")
            sec_title = sec.get("title", "")
            if sec_title and len(sec_title) > 4 and not re.search(re.escape(sec_title), text, re.IGNORECASE):
                issues.append({
                    "severity": "p0",
                    "check": "section_presence",
                    "message": f"章节标题缺失：'{sec_title}'",
                    "section_id": sec_id,
                })

        # Check 3: numeric baseline claims — P1-5: granular threshold
        numeric_baseline = self.ctx.understanding.get("numeric_baseline", {})
        intent = self.ctx.understanding.get("intent", "fresh")
        if numeric_baseline and intent in ("fill_from_reference", "extend"):
            all_evidence = "\n".join(self.ctx.research_findings.values())
            unverified_ratio = verify_section_claims(text, numeric_baseline, all_evidence)
            # P1-5: Any ratio > 30% is a P2 warning; > 50% escalates to P1 error.
            if unverified_ratio > 0.5:
                issues.append({
                    "severity": "p1",
                    "check": "numeric_claims",
                    "message": (
                        f"超过 {unverified_ratio:.0%} 的数字声明无法在参考材料中验证，"
                        "存在较高幻觉风险"
                    ),
                    "section_id": None,
                })
            elif unverified_ratio > 0.30:
                issues.append({
                    "severity": "p2",
                    "check": "numeric_claims",
                    "message": (
                        f"{unverified_ratio:.0%} 的数字声明无法在参考材料中验证，"
                        "建议人工复核关键数据"
                    ),
                    "section_id": None,
                })

        # Check 4: PPTX assertion title quality
        from app.rendering.doc_spec import PptxSpec
        if isinstance(self.ctx.spec, PptxSpec):
            issues.extend(_check_pptx_assertion_titles(self.ctx.spec, text))

        # Check 5: cross-section numeric consistency
        issues.extend(_check_cross_section_numeric_consistency(text))

        # Check 6: PPT cross-slide narrative coherence (duplicate/near-duplicate titles)
        from app.rendering.doc_spec import PptxSpec
        if isinstance(self.ctx.spec, PptxSpec):
            issues.extend(_check_pptx_slide_title_uniqueness(self.ctx.spec))

        # Check 7: PPTX visual density — too much bullet text per slide
        if isinstance(self.ctx.spec, PptxSpec):
            issues.extend(_check_pptx_visual_density(
                self.ctx.spec, rendered_bytes=self.ctx.rendered_bytes
            ))

        # Check 8: Detect render-failure placeholder strings left in document
        issues.extend(_check_render_failure_placeholders(text))

        # Check 9: P1-5 Per-section character count vs target_chars (DOCX)
        # Check 9b: P1-2 Per-slide content density check (PPTX)
        from app.rendering.doc_spec import DocxSpec, PptxSpec
        if isinstance(self.ctx.spec, DocxSpec):
            issues.extend(_check_per_section_chars(self.ctx.spec, text))
        if isinstance(self.ctx.spec, PptxSpec):
            issues.extend(_check_pptx_per_slide_content(self.ctx.spec))

        # Check 10: Spec→render equivalence — tables and charts in spec must appear rendered
        issues.extend(_check_spec_render_equivalence(
            self.ctx.spec, rendered_bytes, self.ctx.rendered_ext or "docx"
        ))

        # Check 11: P3-4 Named entity consistency — only check entities in entity_registry
        # (whitelist mode) to avoid false positives from regex-extracted org names.
        entity_registry = self.ctx.understanding.get("entity_registry", {}) if self.ctx.understanding else {}
        issues.extend(_check_named_entity_consistency(text, entity_registry=entity_registry))

        # Check 12: P3-2 DOCX cross-reference numbers — "图N/表N" citations must not
        # exceed the actual count of figures/tables in the rendered document
        if isinstance(self.ctx.spec, DocxSpec):
            issues.extend(_check_cross_reference_numbers(self.ctx.spec, text))

        # Check 13: P3-2 PPTX layout diversity — warn if >80% slides use same layout
        if isinstance(self.ctx.spec, PptxSpec):
            issues.extend(_check_pptx_layout_diversity(self.ctx.spec))

        # Check 14: P2-4 PPTX slide content density — flag over-stuffed slides
        if isinstance(self.ctx.spec, PptxSpec):
            issues.extend(_check_pptx_slide_density(self.ctx.spec))

        # Check 15: P2-5 XLSX cross-sheet numeric consistency
        from app.rendering.doc_spec import XlsxSpec
        if isinstance(self.ctx.spec, XlsxSpec) and rendered_bytes:
            issues.extend(_check_xlsx_cross_sheet_consistency(rendered_bytes))

        # Check 16: P1-2/P3-4 Lightweight hallucination detection — key numbers in
        # rendered text OR PPTX spec text must be traceable to research_findings.
        # PPTX: extract numbers from assertion_titles + bullets (rendered_bytes may
        # be unavailable for pptx in the first QA pass).
        if self.ctx.research_findings:
            all_evidence = "\n".join(self.ctx.research_findings.values())
            check_text = text
            if isinstance(self.ctx.spec, PptxSpec) and not check_text:
                # Build text corpus from PPTX spec directly
                pptx_lines: list[str] = []
                for sl in (self.ctx.spec.slides or []):
                    pptx_lines.append(sl.assertion_title)
                    pptx_lines.extend(sl.bullets or [])
                    if sl.speaker_notes:
                        pptx_lines.append(sl.speaker_notes)
                check_text = "\n".join(pptx_lines)
            if check_text:
                issues.extend(_check_hallucination_numbers(check_text, all_evidence))

        # Check 17: P2-3 Cross-section duplicate paragraph detection
        # Finds paragraphs repeated verbatim or near-verbatim across sections,
        # which indicates copy-paste failures in LLM generation.
        from app.rendering.doc_spec import DocxSpec
        if isinstance(self.ctx.spec, (DocxSpec, PptxSpec)):
            issues.extend(_check_cross_section_duplicate_paragraphs(self.ctx.spec))

        # Check 18: P2-4 PPTX section_header slide ordering — divider slides should
        # immediately precede a group of content slides, never appear consecutively or
        # at the end.  Catches outlines where the LLM stacked multiple section_header
        # slides without interleaving content.
        if isinstance(self.ctx.spec, PptxSpec):
            issues.extend(_check_section_header_ordering(self.ctx.spec))

        # Check 19: P2-4 PPTX big_number speaker_notes number coherence.
        # big_number slides present a headline KPI — the speaker_notes should
        # reference the same number so the presenter can explain it.
        if isinstance(self.ctx.spec, PptxSpec):
            issues.extend(_check_big_number_notes_coherence(self.ctx.spec))

        # Check 20: P3-1 PPTX final slide should be a conclusion/CTA slide.
        # Professional decks end with a "next steps", "conclusion" or "Q&A" slide;
        # a deck that ends mid-content feels unfinished and unpresentable.
        if isinstance(self.ctx.spec, PptxSpec):
            issues.extend(_check_pptx_final_slide_cta(self.ctx.spec))

        # Check 21: R16-P1 DOCX non-last sections should end with a transition sentence.
        # The spec_gen system prompt requires this, but LLMs sometimes omit it.
        # Catching it here allows a targeted QA re-run for the offending section.
        if isinstance(self.ctx.spec, DocxSpec):
            issues.extend(_check_docx_section_transitions(self.ctx.spec))

        # Check 22: R16-P1 PPTX speaker_notes should be substantive (≥60 chars).
        # Notes shorter than 60 chars are usually placeholder text and give
        # presenters insufficient material to explain the slide.
        if isinstance(self.ctx.spec, PptxSpec):
            issues.extend(_check_pptx_speaker_notes_length(self.ctx.spec))

        # Check 23: R16-P2 PPTX agenda slide bullets should match actual slide titles.
        # Auto-generated agenda slides can drift from actual content if titles were
        # edited during critique or QA passes.
        if isinstance(self.ctx.spec, PptxSpec):
            issues.extend(_check_pptx_agenda_sync(self.ctx.spec))

        # Check 24: R17-P1 Excel cross-sheet formula reference validation.
        # Formulas like ='SheetName'!A1 that reference nonexistent sheets produce
        # silent #REF! errors — catch them before the user opens the file.
        from app.rendering.doc_spec import XlsxSpec
        if isinstance(self.ctx.spec, XlsxSpec) and rendered_bytes:
            issues.extend(_check_xlsx_formula_cross_refs(rendered_bytes))

        # Check 25: R17-P2 DOCX section word count drop-off detection.
        # LLMs sometimes lose coherence toward the end of long documents, generating
        # increasingly thin content in the final sections.
        if isinstance(self.ctx.spec, DocxSpec) and text:
            issues.extend(_check_docx_section_char_dropoff(self.ctx.spec, text))

        # Check 26: R17-P2 fill_from_reference sections should cite evidence sources.
        # Sections with empty evidence_ids in a reference-based document are likely
        # hallucinated without grounding, since the research pipeline provides material.
        intent = self.ctx.understanding.get("intent", "fresh") if self.ctx.understanding else "fresh"
        has_research = bool(self.ctx.research_findings)
        if isinstance(self.ctx.spec, DocxSpec):
            issues.extend(_check_docx_evidence_coverage(self.ctx.spec, intent, has_research))

        # Check 27: R18-P2 DOCX last section must end with a closing/summary sentence.
        # Complement to Check 21 (transition sentences for non-last sections).
        # The final section should conclude the document, not end abruptly mid-topic.
        if isinstance(self.ctx.spec, DocxSpec) and text:
            issues.extend(_check_docx_last_section_closing(self.ctx.spec, text))

        # Check 28: R18-P1 Excel named range conflict detection in rendered workbook.
        # Named ranges that shadow or conflict with sheet-level names cause subtle
        # formula errors that are invisible until the user edits the workbook.
        from app.rendering.doc_spec import XlsxSpec
        if isinstance(self.ctx.spec, XlsxSpec) and rendered_bytes:
            issues.extend(_check_xlsx_named_range_conflicts(rendered_bytes))

        return issues


def _extract_text_from_bytes(rendered_bytes: bytes, ext: str) -> str:
    """Extract plain text from rendered document bytes."""
    import io
    try:
        if ext == "docx":
            from docx import Document
            doc = Document(io.BytesIO(rendered_bytes))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        elif ext == "pptx":
            from pptx import Presentation
            prs = Presentation(io.BytesIO(rendered_bytes))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        parts.append(shape.text_frame.text)
            return "\n".join(parts)
        elif ext == "xlsx":
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(rendered_bytes))
            parts = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    parts.append(" ".join(str(c or "") for c in row))
            return "\n".join(parts)
    except Exception as exc:
        logger.warning("[QA] Text extraction failed for %s: %s", ext, exc)
    return ""


def _iter_sections(spec) -> list[dict]:
    """Yield section dicts from a DocumentSpec."""
    from app.rendering.doc_spec import DocxSpec, PptxSpec, XlsxSpec
    if isinstance(spec, DocxSpec):
        return [{"id": s.id, "title": s.title} for s in spec.sections]
    elif isinstance(spec, PptxSpec):
        return [{"id": s.id, "title": s.assertion_title} for s in spec.slides]
    elif isinstance(spec, XlsxSpec):
        return [{"id": s.id, "title": s.name} for s in spec.sheets]
    return []


def _sum_target_chars(spec) -> int:
    """Sum all section/slide target_chars from a DocumentSpec."""
    from app.rendering.doc_spec import DocxSpec, PptxSpec, XlsxSpec
    if isinstance(spec, DocxSpec):
        return sum(getattr(s, "target_chars", 400) for s in spec.sections)
    if isinstance(spec, PptxSpec):
        return len(spec.slides) * 200  # ~200 chars per PPT slide as proxy
    if isinstance(spec, XlsxSpec):
        return len(spec.sheets) * 300
    return 0


def _check_pptx_assertion_titles(spec, rendered_text: str) -> list[dict]:
    """P1: PPT slide titles must contain a number or clear judgment word."""
    import re
    issues: list[dict] = []
    _JUDGMENT_WORDS = {"增长", "下降", "超过", "达到", "突破", "低于", "完成", "实现", "领先", "落后",
                       "优于", "差于", "最高", "最低", "提升", "下滑", "回升", "持续", "首次", "创新高"}
    number_pattern = re.compile(r"\d")

    for slide in spec.slides:
        title = slide.assertion_title
        has_number = bool(number_pattern.search(title))
        has_judgment = any(w in title for w in _JUDGMENT_WORDS)
        if not has_number and not has_judgment:
            issues.append({
                "severity": "p1",
                "check": "assertion_title_quality",
                "message": f"PPT幻灯片标题缺乏断言性（无数字或判断词）：'{title}'",
                "section_id": slide.id,
            })
    return issues


def _check_pptx_slide_title_uniqueness(spec) -> list[dict]:
    """P2: Detect near-duplicate slide titles which indicate copy-paste hallucination."""
    issues: list[dict] = []
    seen: dict[str, str] = {}  # normalised_title → first slide id

    for slide in spec.slides:
        # Normalise: strip numbers and punctuation for comparison
        normalised = re.sub(r"[\d%，,。！？·\s]", "", slide.assertion_title)
        if len(normalised) < 4:
            continue
        if normalised in seen:
            issues.append({
                "severity": "p1",
                "check": "slide_title_uniqueness",
                "message": (
                    f"幻灯片标题与页 '{seen[normalised]}' 高度相似，"
                    f"可能重复：'{slide.assertion_title}'"
                ),
                "section_id": slide.id,
            })
        else:
            seen[normalised] = slide.id
    return issues


def _check_pptx_visual_density(spec, rendered_bytes: bytes | None = None) -> list[dict]:
    """P1-6/P2-A: Flag slides where bullet text would overflow the body placeholder.

    When rendered_bytes are available (post-render QA pass), actual placeholder
    heights are read from the PPTX file to compute max row capacity more accurately.
    Falls back to a character-count heuristic when bytes are not yet available.
    """
    issues: list[dict] = []

    # P1-6: Read actual body placeholder heights from rendered bytes (EMU → inches)
    slide_ph_heights: dict[int, float] = {}
    if rendered_bytes:
        try:
            import io as _io
            from pptx import Presentation as _Prs
            _prs = _Prs(_io.BytesIO(rendered_bytes))
            for idx, _slide in enumerate(_prs.slides):
                for ph in _slide.placeholders:
                    if ph.placeholder_format.idx != 0 and ph.height:
                        slide_ph_heights[idx] = ph.height / 914400  # EMU → inches
                        break
        except Exception:
            pass

    for slide_idx, slide in enumerate(spec.slides):
        total_bullet_chars = sum(len(b) for b in slide.bullets)
        if not total_bullet_chars:
            continue

        ph_height_in = slide_ph_heights.get(slide_idx)
        if ph_height_in:
            # ~15 Chinese chars per line at 18pt; ~0.35 inches per line
            estimated_rows = total_bullet_chars / 15
            max_rows = ph_height_in / 0.35
            if estimated_rows > max_rows:
                issues.append({
                    "severity": "p1",
                    "check": "pptx_visual_density",
                    "message": (
                        f"幻灯片 '{slide.assertion_title}' 内容可能溢出占位框 "
                        f"（估算 {estimated_rows:.0f} 行，框高约 {max_rows:.0f} 行），"
                        "请压缩要点或拆分为多页"
                    ),
                    "section_id": slide.id,
                })
        elif total_bullet_chars > _SLIDE_BODY_CHAR_CAPACITY:
            # Heuristic fallback (pre-render or extraction failed)
            issues.append({
                "severity": "p1",
                "check": "pptx_visual_density",
                "message": (
                    f"幻灯片 '{slide.assertion_title}' 文字量过大 "
                    f"（{total_bullet_chars} 字符，建议 ≤{_SLIDE_BODY_CHAR_CAPACITY}），"
                    "请压缩内容或拆分为多页"
                ),
                "section_id": slide.id,
            })

    return issues


def _check_render_failure_placeholders(text: str) -> list[dict]:
    """P1-E / P1-2: Flag documents that contain render-failure placeholder strings.

    Detects two classes of failure placeholders:
    1. Chart/table render failures: "[图表：XXX（渲染失败，请手动插入）]"
    2. DOCX section content failures: "[本章节「...」内容生成失败，请重新生成]"
       (inserted by spec_gen's placeholder degradation path when all retries fail)
    """
    issues: list[dict] = []

    # Class 1: chart/table render failures (p1)
    chart_pat = re.compile(r"\[图表[：:].*?渲染失败.*?\]")
    chart_matches = chart_pat.findall(text)
    if chart_matches:
        sample = chart_matches[0][:60]
        issues.append({
            "severity": "p1",
            "check": "render_failure_placeholder",
            "message": (
                f"文档包含 {len(chart_matches)} 处图表渲染失败占位符（例：'{sample}'），"
                "请检查图表数据或手动替换"
            ),
            "section_id": None,
        })

    # Class 2: DOCX section content generation failures (p0 — whole section missing)
    section_pat = re.compile(r"\[本章节[「『](.{1,30})[」』]内容生成失败")
    section_matches = section_pat.findall(text)
    if section_matches:
        sample_title = section_matches[0][:30]
        issues.append({
            "severity": "p0",
            "check": "section_content_generation_failure",
            "message": (
                f"文档包含 {len(section_matches)} 个章节内容生成失败（例：'{sample_title}'），"
                "这些章节需要重新生成"
            ),
            "section_id": None,
        })

    return issues


def _collect_sections_flat(sections) -> list:
    """P2-5: Recursively collect all sections and subsections into a flat list."""
    result = []
    for sec in (sections or []):
        result.append(sec)
        subsections = getattr(sec, "subsections", None) or []
        result.extend(_collect_sections_flat(subsections))
    return result


def _check_per_section_chars(spec, text: str) -> list[dict]:
    """P1-5: Compare actual rendered chars per section against target_chars.

    Splits rendered text at section headings, counts chars in each section body,
    and flags sections whose actual content is under 40% of target_chars.

    P2-5: Also checks subsections recursively, not just top-level sections.
    """
    issues: list[dict] = []
    # P2-5: Flatten top-level sections + subsections
    all_sections = _collect_sections_flat(getattr(spec, "sections", []))
    if not all_sections or not text:
        return issues

    # Build title → (section_id, target_chars) map
    title_map = [
        (s.title, s.id, s.target_chars)
        for s in all_sections
        if s.title and getattr(s, "target_chars", 0) > 100
    ]
    if not title_map:
        return issues

    for i, (title, sec_id, target) in enumerate(title_map):
        title_pos = text.find(title)
        if title_pos < 0:
            continue  # section_presence check (Check 2) handles missing titles

        # Find where the next section starts
        next_pos = len(text)
        for j, (other_title, _, _) in enumerate(title_map):
            if j == i:
                continue
            other_pos = text.find(other_title, title_pos + len(title))
            if 0 < other_pos < next_pos:
                next_pos = other_pos

        section_body = text[title_pos + len(title):next_pos]
        actual_chars = len(re.sub(r"\s+", "", section_body))

        if actual_chars < target * 0.40:
            issues.append({
                "severity": "p1",
                "check": "per_section_chars",
                "message": (
                    f"章节「{title}」字数不足（渲染后约 {actual_chars} 字，"
                    f"目标 {target} 字，低于 40% 阈值），可能生成失败或被截断"
                ),
                "section_id": sec_id,
            })

    return issues


def _check_cross_section_numeric_consistency(text: str) -> list[dict]:
    """P2: Detect same-label numbers that differ across sections (probable hallucination)."""
    import re
    issues: list[dict] = []

    # Find patterns like "XX率 N%" or "XX额 N万" with their surrounding label
    pattern = re.compile(r"([一-龥]{2,8})[：:为达到约]?\s*(\d[\d,\.]+)\s*(%|万|亿|元|人|个|次)?")
    findings: dict[str, list[str]] = {}
    for m in pattern.finditer(text):
        label = m.group(1)
        value = m.group(2).replace(",", "")
        unit = m.group(3) or ""
        key = f"{label}{unit}"
        findings.setdefault(key, [])
        findings[key].append(value)

    for key, values in findings.items():
        unique = set(values)
        if len(unique) > 1 and len(values) >= 2:
            # Only flag when values differ by more than 5% (avoid rounding noise)
            try:
                nums = [float(v) for v in unique]
                mn, mx = min(nums), max(nums)
                if mn > 0 and (mx - mn) / mn > 0.05:
                    issues.append({
                        "severity": "p1",
                        "check": "cross_section_consistency",
                        "message": f"'{key}' 在文档中出现不一致数值：{sorted(unique)}",
                        "section_id": None,
                    })
            except ValueError:
                pass

    return issues


def _check_spec_render_equivalence(spec, rendered_bytes: bytes, ext: str) -> list[dict]:
    """Check 10: Tables and charts declared in spec actually appear in the rendered document.

    Uses python-docx / python-pptx / openpyxl to count structural elements in the
    rendered bytes and compares against spec declarations.  Mismatches indicate a
    silent renderer failure (e.g. bad data shape, missing dependency).
    """
    import io
    issues: list[dict] = []

    try:
        from app.rendering.doc_spec import DocxSpec, PptxSpec, XlsxSpec

        if isinstance(spec, DocxSpec) and ext == "docx":
            try:
                from docx import Document
                doc = Document(io.BytesIO(rendered_bytes))
                rendered_tables = len(doc.tables)
                rendered_pics = sum(
                    1
                    for p in doc.paragraphs
                    for r in p.runs
                    if r._r.findall(
                        ".//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}drawing"
                    )
                )
            except Exception:
                return issues

            spec_tables = sum(1 for s in spec.sections if s.table)
            spec_charts = sum(1 for s in spec.sections if s.chart)

            if spec_tables > 0 and rendered_tables < spec_tables:
                issues.append({
                    "severity": "p1",
                    "check": "spec_render_equivalence",
                    "message": (
                        f"规格中声明了 {spec_tables} 个表格，"
                        f"但渲染文档中只检测到 {rendered_tables} 个，"
                        "部分表格可能未成功渲染"
                    ),
                    "section_id": None,
                })
            if spec_charts > 0 and rendered_pics < spec_charts:
                issues.append({
                    "severity": "p1",
                    "check": "spec_render_equivalence",
                    "message": (
                        f"规格中声明了 {spec_charts} 个图表，"
                        f"但渲染文档中只检测到 {rendered_pics} 张图片，"
                        "部分图表可能渲染失败"
                    ),
                    "section_id": None,
                })

        elif isinstance(spec, PptxSpec) and ext == "pptx":
            try:
                from pptx import Presentation
                prs = Presentation(io.BytesIO(rendered_bytes))
                rendered_slide_count = len(prs.slides)
            except Exception:
                return issues

            spec_slide_count = len(spec.slides)
            if rendered_slide_count < spec_slide_count:
                issues.append({
                    "severity": "p0",
                    "check": "spec_render_equivalence",
                    "message": (
                        f"规格中声明了 {spec_slide_count} 张幻灯片，"
                        f"但渲染文件只有 {rendered_slide_count} 张，"
                        "存在幻灯片丢失"
                    ),
                    "section_id": None,
                })

        elif isinstance(spec, XlsxSpec) and ext == "xlsx":
            try:
                import openpyxl
                wb = openpyxl.load_workbook(io.BytesIO(rendered_bytes))
                rendered_sheet_count = len(wb.worksheets)
            except Exception:
                return issues

            spec_sheet_count = len(spec.sheets)
            # Rendered workbook includes a dashboard sheet (+1) not in spec
            if rendered_sheet_count < spec_sheet_count:
                issues.append({
                    "severity": "p1",
                    "check": "spec_render_equivalence",
                    "message": (
                        f"规格中声明了 {spec_sheet_count} 个工作表，"
                        f"但渲染文件只有 {rendered_sheet_count} 个，"
                        "部分工作表可能生成失败"
                    ),
                    "section_id": None,
                })

    except Exception as exc:
        logger.debug("[QA] spec_render_equivalence check error (non-fatal): %s", exc)

    return issues


def _check_named_entity_consistency(
    text: str,
    entity_registry: dict | None = None,
) -> list[dict]:
    """Check 11: P3-4 Detect named entity inconsistency using whitelist mode when available.

    Whitelist mode (entity_registry provided and non-empty):
      Only check canonical entity names from entity_registry. For each canonical name,
      look for alternate forms (base name without suffix, or other suffixes) that appear
      in the text — flag when the canonical and an alternate co-exist.

    Fallback mode (no entity_registry):
      Use the original regex-based heuristic that extracts all org-like names and
      clusters them by base name (may produce false positives).
    """
    issues: list[dict] = []
    try:
        _ORG_SUFFIXES = (
            "公司", "集团", "部门", "中心", "研究院", "研究所", "事业部",
            "子公司", "分公司", "单位", "院", "所", "局", "厅",
        )

        if entity_registry:
            # Whitelist mode — only check entities we explicitly know about
            for canonical, _meta in entity_registry.items():
                if not canonical or len(canonical) < 3:
                    continue
                # Derive base name by stripping suffix
                base = canonical
                for suf in _ORG_SUFFIXES:
                    if canonical.endswith(suf) and len(canonical) > len(suf) + 1:
                        base = canonical[: -len(suf)]
                        break
                # Find alternate forms: base + other suffixes that appear in text
                alternates = set()
                for suf in _ORG_SUFFIXES:
                    alt = base + suf
                    if alt != canonical and alt in text:
                        alternates.add(alt)
                if alternates and canonical in text:
                    sample = "、".join(sorted(alternates)[:3])
                    issues.append({
                        "severity": "p1",
                        "check": "named_entity_consistency",
                        "message": (
                            f"实体「{canonical}」与变体称谓「{sample}」在文档中同时出现，"
                            "请统一称谓"
                        ),
                        "section_id": None,
                    })
        else:
            # Fallback: regex-based heuristic (may produce false positives)
            org_pattern = re.compile(
                r"[^\s，。、；：！？（）【】“”‘’]{2,12}"
                r"(?:公司|集团|部门|中心|研究院|研究所|事业部|子公司|分公司|单位|院|所|局|厅)"
            )
            found: dict[str, set] = {}
            for m in org_pattern.finditer(text):
                full = m.group()
                base = full
                for suf in _ORG_SUFFIXES:
                    if full.endswith(suf) and len(full) > len(suf) + 1:
                        base = full[: -len(suf)]
                        break
                if len(base) < 2:
                    continue
                found.setdefault(base, set()).add(full)

            for base, names in found.items():
                if len(names) >= 2:
                    sample = "、".join(sorted(names)[:4])
                    issues.append({
                        "severity": "p1",
                        "check": "named_entity_consistency",
                        "message": (
                            f"实体「{base}」在文档中以 {len(names)} 种不同名称出现："
                            f"{sample}，请统一称谓"
                        ),
                        "section_id": None,
                    })

    except Exception as exc:
        logger.debug("[QA] named_entity_consistency check error (non-fatal): %s", exc)
    return issues


def _check_cross_reference_numbers(spec, text: str) -> list[dict]:
    """Check 12: P3-2 Validate that 'Figure N' / 'Table N' citations in body text do not
    exceed the actual count of figures and tables in the spec.

    When text says '如表5所示' but the spec only has 3 tables, the citation is wrong
    (probably a copy-paste error or hallucinated reference).
    """
    issues: list[dict] = []
    try:
        spec_table_count = sum(1 for s in spec.sections if s.table)
        spec_chart_count = sum(1 for s in spec.sections if s.chart)

        # Find highest cited table / figure number in body text
        table_refs = re.findall(r"表\s*(\d+)", text)
        chart_refs = re.findall(r"图\s*(\d+)", text)

        if table_refs:
            max_cited = max(int(n) for n in table_refs)
            if spec_table_count > 0 and max_cited > spec_table_count:
                issues.append({
                    "severity": "p1",
                    "check": "cross_reference_numbers",
                    "message": (
                        f"正文引用了表{max_cited}，但规格中只有 {spec_table_count} 个表格，"
                        "交叉引用编号可能错误"
                    ),
                    "section_id": None,
                })

        if chart_refs:
            max_cited = max(int(n) for n in chart_refs)
            if spec_chart_count > 0 and max_cited > spec_chart_count:
                issues.append({
                    "severity": "p1",
                    "check": "cross_reference_numbers",
                    "message": (
                        f"正文引用了图{max_cited}，但规格中只有 {spec_chart_count} 个图表，"
                        "交叉引用编号可能错误"
                    ),
                    "section_id": None,
                })
    except Exception as exc:
        logger.debug("[QA] cross_reference_numbers check error (non-fatal): %s", exc)
    return issues


def _check_pptx_layout_diversity(spec) -> list[dict]:
    """Check 13: P3-2 Flag when >80% of slides use the same layout.

    Monotonous layout use (e.g. all 'content') produces visually dull decks.
    This is a warning (p2) that suggests the LLM did not vary layouts enough.
    """
    issues: list[dict] = []
    try:
        slides = getattr(spec, "slides", [])
        if len(slides) < 4:
            return issues  # Too few slides to bother

        from collections import Counter
        layout_counts = Counter(s.layout for s in slides)
        most_common_layout, most_common_count = layout_counts.most_common(1)[0]
        ratio = most_common_count / len(slides)

        if ratio > 0.80:
            issues.append({
                "severity": "p2",
                "check": "pptx_layout_diversity",
                "message": (
                    f"{most_common_count}/{len(slides)} 张幻灯片使用相同布局 '{most_common_layout}'，"
                    "建议增加 comparison / big_number / section_header 等布局以提升视觉多样性"
                ),
                "section_id": None,
            })
    except Exception as exc:
        logger.debug("[QA] pptx_layout_diversity check error (non-fatal): %s", exc)
    return issues


def _check_pptx_slide_density(spec) -> list[dict]:
    """Check 14: P2-4 Flag slides whose bullet text total exceeds layout-specific limits.

    Over-stuffed slides (>300 chars in content layout) are unreadable in presentations.
    The renderer already truncates at these limits (P1-2), so this QA check surfaces
    cases where the spec itself was over-specified before rendering truncation.
    """
    issues: list[dict] = []
    _LIMITS = {"content": 300, "comparison": 200, "big_number": 80, "section_header": 60}
    try:
        for slide in getattr(spec, "slides", []):
            total = sum(len(b) for b in (slide.bullets or []))
            limit = _LIMITS.get(slide.layout, 300)
            if total > limit:
                issues.append({
                    "severity": "p2",
                    "check": "pptx_slide_density",
                    "message": (
                        f"幻灯片「{slide.assertion_title[:30]}」（{slide.layout}）"
                        f"bullet 总字数 {total} 超过布局上限 {limit}，"
                        "已自动截断，建议精简内容"
                    ),
                    "section_id": slide.id,
                })
    except Exception as exc:
        logger.debug("[QA] pptx_slide_density check error (non-fatal): %s", exc)
    return issues


def _check_xlsx_cross_sheet_consistency(rendered_bytes: bytes) -> list[dict]:
    """Check 15: P2-5 Detect same-label metrics with different values across XLSX sheets.

    Extracts numeric values associated with common metric labels from each sheet,
    then flags label-value pairs that differ across sheets with the same label.
    """
    issues: list[dict] = []
    try:
        import io
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(rendered_bytes), read_only=True, data_only=True)

        # label → {sheet_name: value_str}
        label_sheet_values: dict[str, dict[str, str]] = {}

        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                cells = [c for c in row if c is not None]
                # Look for (label, numeric_value) pairs in adjacent cells
                for i in range(len(cells) - 1):
                    label = str(cells[i]).strip()
                    val = cells[i + 1]
                    if (
                        len(label) >= 2
                        and isinstance(val, (int, float))
                        and label not in ("", "None")
                    ):
                        label_sheet_values.setdefault(label, {})[ws.title] = str(val)

        for label, sheet_vals in label_sheet_values.items():
            if len(sheet_vals) >= 2:
                unique_vals = set(sheet_vals.values())
                if len(unique_vals) > 1:
                    detail = "、".join(f"{s}={v}" for s, v in list(sheet_vals.items())[:3])
                    issues.append({
                        "severity": "p2",
                        "check": "xlsx_cross_sheet_consistency",
                        "message": (
                            f"指标「{label}」在多个Sheet中数值不一致：{detail}，"
                            "请确认数据来源或使用跨表引用公式"
                        ),
                        "section_id": None,
                    })
    except Exception as exc:
        logger.debug("[QA] xlsx_cross_sheet_consistency check error (non-fatal): %s", exc)
    return issues


def _normalize_number_str(s: str) -> str:
    """P2-5: Normalize numeric strings to a canonical form for evidence matching.

    Handles:
    - Chinese unit suffixes: 亿 (×1e8), 万 (×1e4), 千 (×1e3) → plain integer
    - Thousand-separator commas: '1,234,567' → '1234567'
    - Trailing % preserved after digit normalization

    This prevents false positives when evidence says '1.5亿' and the document
    says '150,000,000' — both represent the same value.
    """
    s = s.strip()
    is_pct = s.endswith("%")
    core = s.rstrip("%").replace(",", "")  # strip commas and trailing %
    try:
        val = float(core)
    except ValueError:
        return s  # cannot parse — return as-is
    if "亿" in s:
        val *= 1e8
    elif "万" in s and "亿" not in s:
        val *= 1e4
    elif "千" in s and "万" not in s:
        val *= 1e3
    # Represent as integer when no fractional part, else float with up to 2 dp
    if val == int(val):
        normalized = str(int(val))
    else:
        normalized = f"{val:.2f}".rstrip("0").rstrip(".")
    return normalized + "%" if is_pct else normalized


def _check_hallucination_numbers(text: str, all_evidence: str) -> list[dict]:
    """Check 16: P3-4 Lightweight hallucination detection for numeric claims.

    Extracts 4+-digit numbers and percentage figures from the rendered document,
    then checks whether each appears in the research evidence. Numbers not found
    in any evidence block are potential hallucinations.

    P2-5: Normalizes 亿/万/千分位 before matching so '1.5亿' and '150000000'
    are treated as equivalent, reducing false positives.

    Only fires when evidence is non-empty and at least 3 suspicious numbers found.
    Conservative threshold: flags only when >40% of checked numbers are ungrounded.
    """
    issues: list[dict] = []
    try:
        if not all_evidence or not text:
            return issues

        # Extract significant numbers: 4+ digit integers, and percentages
        num_pattern = re.compile(r"\b\d{4,}\b|\b\d+\.\d+%|\b\d+%")
        doc_numbers = set(num_pattern.findall(text))
        if len(doc_numbers) < 3:
            return issues  # too few numbers to make a meaningful judgment

        # P2-5: Build a normalized evidence corpus for matching
        evidence_normalized = set()
        for ev_num in num_pattern.findall(all_evidence):
            evidence_normalized.add(ev_num)
            evidence_normalized.add(_normalize_number_str(ev_num))

        ungrounded = []
        for n in doc_numbers:
            n_norm = _normalize_number_str(n)
            if n not in all_evidence and n_norm not in evidence_normalized:
                ungrounded.append(n)

        ratio = len(ungrounded) / len(doc_numbers)

        if ratio > 0.40 and len(ungrounded) >= 3:
            sample = "、".join(list(ungrounded)[:5])
            issues.append({
                "severity": "p2",
                "check": "hallucination_numbers",
                "message": (
                    f"{len(ungrounded)}/{len(doc_numbers)} 个数字（{sample}…）"
                    "在研究材料中未找到依据，可能存在幻觉，请人工核实"
                ),
                "section_id": None,
            })
    except Exception as exc:
        logger.debug("[QA] hallucination_numbers check error (non-fatal): %s", exc)
    return issues


def _check_cross_section_duplicate_paragraphs(spec) -> list[dict]:
    """Check 17: P3-3 Detect near-duplicated content across DOCX sections and PPTX slides.

    - DOCX: checks paragraphs across sections (≥30 chars), threshold 0.65
    - PPTX: checks bullet points across slides (≥20 chars), threshold 0.65

    Threshold lowered from 0.80 to 0.65 to catch semantically similar rewrites
    (same sentence with minor paraphrasing) which the LLM commonly produces when
    instructed to "expand on" a topic it already covered.
    """
    _THRESHOLD = 0.65
    issues: list[dict] = []

    def _ngrams4(text: str) -> set:
        return {text[i:i+4] for i in range(len(text) - 3)} if len(text) >= 4 else {text}

    def _jaccard(a: str, b: str) -> float:
        sa, sb = _ngrams4(a), _ngrams4(b)
        if not sa and not sb:
            return 1.0
        union = len(sa | sb)
        return len(sa & sb) / union if union else 0.0

    def _scan_texts(items: list[tuple[str, str, str]], min_len: int, check_name: str) -> None:
        """items = list of (text, container_id, container_title)."""
        seen: list[tuple[str, str, str]] = []
        for text_a, cid_a, ctitle_a in items:
            for text_b, cid_b, ctitle_b in seen:
                if cid_a == cid_b:
                    continue
                sim = _jaccard(text_a, text_b)
                if sim >= _THRESHOLD:
                    snippet = text_a[:50].replace("\n", " ")
                    issues.append({
                        "severity": "p1",
                        "check": check_name,
                        "message": (
                            f"「{ctitle_a}」与「{ctitle_b}」含近似重复内容"
                            f"（相似度 {sim:.0%}）：\"{snippet}...\"，请改写以避免重复"
                        ),
                        "section_id": cid_a,
                    })
                    break
            seen.append((text_a, cid_a, ctitle_a))

    try:
        # DOCX: check paragraphs across sections
        all_sections = _collect_sections_flat(getattr(spec, "sections", []))
        if all_sections:
            paras: list[tuple[str, str, str]] = []
            for sec in all_sections:
                for p in (sec.paragraphs or []):
                    if p and len(p.strip()) >= 30:
                        paras.append((p.strip(), sec.id, sec.title))
            if len(paras) >= 2:
                _scan_texts(paras, 30, "cross_section_duplicate_paragraph")

        # PPTX: check bullet points across slides
        slides = getattr(spec, "slides", [])
        if slides:
            bullets: list[tuple[str, str, str]] = []
            for slide in slides:
                slide_label = getattr(slide, "assertion_title", getattr(slide, "id", ""))
                slide_id = getattr(slide, "id", "")
                for b in (getattr(slide, "bullets", []) or []):
                    if b and len(b.strip()) >= 20:
                        bullets.append((b.strip(), slide_id, slide_label))
            if len(bullets) >= 2:
                _scan_texts(bullets, 20, "cross_slide_duplicate_bullet")

    except Exception as exc:
        logger.debug("[QA] cross_section_duplicate check error (non-fatal): %s", exc)
    return issues


def _check_section_header_ordering(spec) -> list[dict]:
    """Check 18: P2-4 PPTX section_header slides must precede content, not stack or trail.

    Valid pattern:   section_header → content → content → section_header → content
    Invalid patterns:
    - Two consecutive section_header slides (stacking)
    - A section_header slide at the very end of the deck (dangling divider)
    """
    issues: list[dict] = []
    try:
        slides = list(spec.slides or [])
        if len(slides) < 3:
            return issues

        layouts = [getattr(s, "layout", "") or "" for s in slides]
        stacked: list[str] = []

        for i in range(len(layouts) - 1):
            if layouts[i] == "section_header" and layouts[i + 1] == "section_header":
                # P3-3: Skip the first pair — cover + first section_header is a valid
                # and common PPTX opening pattern (title slide → agenda/first chapter).
                if i == 0:
                    continue
                stacked.append(f"幻灯片 {i + 1} & {i + 2}（'{slides[i].assertion_title[:20]}' / '{slides[i + 1].assertion_title[:20]}'）")

        if stacked:
            issues.append({
                "severity": "p1",
                "check": "section_header_stacking",
                "message": (
                    f"连续 section_header 幻灯片（{len(stacked)} 处）：{stacked[0]}；"
                    "章节分隔页之间应有内容页，请在两个分隔页之间插入内容幻灯片"
                ),
                "section_id": None,
            })

        # Check for dangling section_header at the end
        if layouts and layouts[-1] == "section_header":
            issues.append({
                "severity": "p1",
                "check": "section_header_trailing",
                "message": (
                    f"最后一张幻灯片是 section_header（'{slides[-1].assertion_title[:30]}'），"
                    "章节分隔页后必须有内容页，否则为空章节"
                ),
                "section_id": getattr(slides[-1], "id", None),
            })

    except Exception as exc:
        logger.debug("[QA] section_header ordering check error (non-fatal): %s", exc)
    return issues


def _check_pptx_per_slide_content(spec) -> list[dict]:
    """Check 9b: P1-2 Flag PPTX slides with no substantive content.

    A slide must have at least one of:
    - assertion_title with ≥5 characters
    - at least 1 bullet point with ≥10 characters
    - a table or chart spec

    Slides with layout='section_header' are exempt (they are intentionally minimal).
    """
    issues: list[dict] = []
    try:
        for slide in (spec.slides or []):
            layout = getattr(slide, "layout", "") or ""
            if layout == "section_header":
                continue  # divider slides are intentionally sparse

            title = (getattr(slide, "assertion_title", "") or "").strip()
            bullets = [b for b in (getattr(slide, "bullets", []) or []) if b and len(b.strip()) >= 10]
            has_table = getattr(slide, "table", None) is not None
            has_chart = getattr(slide, "chart", None) is not None

            if len(title) < 5 and not bullets and not has_table and not has_chart:
                # P2-2: big_number empty is p0 — it produces a completely blank KPI card
                # which is worse than a placeholder (visually broken, not just thin).
                severity = "p0" if layout == "big_number" else "p1"
                issues.append({
                    "severity": severity,
                    "check": "pptx_empty_slide",
                    "message": (
                        f"幻灯片「{getattr(slide, 'id', '?')}」（layout={layout}）内容为空，"
                        "标题不足5字且无要点、表格或图表，请补充内容"
                    ),
                    "section_id": getattr(slide, "id", None),
                })
    except Exception as exc:
        logger.debug("[QA] pptx_per_slide_content check error (non-fatal): %s", exc)
    return issues


def _check_big_number_notes_coherence(spec) -> list[dict]:
    """Check 19: P2-4 big_number slides should reference their headline number in speaker_notes.

    A 'big_number' layout slide presents a single KPI number as the visual focus.
    If the speaker_notes exist but don't mention any digit from the assertion_title,
    the presenter's script and the slide are misaligned.
    """
    issues: list[dict] = []
    try:
        for slide in (spec.slides or []):
            if getattr(slide, "layout", "") != "big_number":
                continue
            notes = (getattr(slide, "speaker_notes", "") or "").strip()
            if not notes:
                continue  # missing notes is caught by speaker_notes sentence check

            title = getattr(slide, "assertion_title", "") or ""
            title_digits = set(re.findall(r"\d+", title))
            if not title_digits:
                continue  # title has no numbers to check

            notes_digits = set(re.findall(r"\d+", notes))
            if not title_digits & notes_digits:
                issues.append({
                    "severity": "p2",
                    "check": "big_number_notes_incoherent",
                    "message": (
                        f"big_number 幻灯片「{title[:30]}」的演讲者备注未提及标题中的关键数字"
                        f"（{', '.join(sorted(title_digits)[:3])}），请在备注中引用该数字加以说明"
                    ),
                    "section_id": getattr(slide, "id", None),
                })
    except Exception as exc:
        logger.debug("[QA] big_number_notes_coherence check error (non-fatal): %s", exc)
    return issues


def _check_pptx_final_slide_cta(spec) -> list[dict]:
    """Check 20: P3-1 PPTX deck should close with a conclusion/CTA/Q&A slide.

    A professional presentation ends with a dedicated closing slide that signals
    completion to the audience. Detecting a missing close-out prevents awkward
    deck endings mid-analysis.
    """
    _CTA_KEYWORDS = {
        "总结", "结论", "建议", "下一步", "行动", "展望", "问答", "Q&A", "Thanks",
        "谢谢", "致谢", "附录", "contact", "conclusion", "next steps", "summary",
        "recommendations", "action", "thank", "end",
    }
    issues: list[dict] = []
    try:
        slides = list(spec.slides or [])
        if len(slides) < 4:
            return issues  # too short to require a dedicated close

        last_title = (getattr(slides[-1], "assertion_title", "") or "").lower()
        has_cta = any(kw.lower() in last_title for kw in _CTA_KEYWORDS)

        if not has_cta:
            issues.append({
                "severity": "p2",
                "check": "pptx_missing_closing_slide",
                "message": (
                    f"演示文稿最后一页「{slides[-1].assertion_title[:30]}」"
                    "不像结束页（应含总结/建议/下一步/Q&A等关键词），"
                    "建议添加一张收尾幻灯片"
                ),
                "section_id": getattr(slides[-1], "id", None),
            })
    except Exception as exc:
        logger.debug("[QA] pptx_final_slide_cta check error (non-fatal): %s", exc)
    return issues


def _check_docx_section_transitions(spec) -> list[dict]:
    """Check 21: R16-P1 Non-last DOCX sections should end with a transition sentence.

    The spec_gen system prompt requires the last paragraph of each non-final section
    to contain a bridging sentence that previews the next chapter.  LLMs occasionally
    ignore this rule.  Catching the failure here enables a targeted QA re-run so the
    section is regenerated with an explicit lookahead prompt (next_section_title).

    Only checks sections with ≥2 paragraphs — short sections may not warrant a
    dedicated transition sentence.
    """
    _TRANSITION_KEYWORDS = (
        "下一", "下部分", "下章", "接下来", "进一步将", "此后", "随后将", "其次",
        "将从", "将在下", "接着", "将深入", "将探讨", "将分析", "将介绍",
    )
    issues: list[dict] = []
    try:
        sections = list(getattr(spec, "sections", []) or [])
        if len(sections) < 2:
            return issues

        for sec in sections[:-1]:  # skip last section — it needs a closing sentence, not a transition
            paras = getattr(sec, "paragraphs", []) or []
            if len(paras) < 2:
                continue  # too short; no expectation of a transition paragraph
            last_para = paras[-1] if paras else ""
            has_transition = any(kw in last_para for kw in _TRANSITION_KEYWORDS)
            if not has_transition and len(last_para) > 20:
                issues.append({
                    "severity": "p2",
                    "check": "docx_section_transition",
                    "message": (
                        f"章节「{sec.title}」末段落缺少承上启下过渡句"
                        "（末段应以「下一部分将…」等语句引出下一章节内容）"
                    ),
                    "section_id": sec.id,
                })
    except Exception as exc:
        logger.debug("[QA] docx_section_transitions check error (non-fatal): %s", exc)
    return issues


def _check_pptx_speaker_notes_length(spec) -> list[dict]:
    """Check 22: R16-P1 PPTX slides should have substantive speaker_notes (≥60 chars).

    Notes shorter than 60 characters are usually placeholder text like '详见幻灯片' or
    a bare repetition of the title.  Presenters need enough notes to speak for 1-2
    minutes per slide — substantive notes enable this.

    section_header slides are exempt since they are intentionally brief.
    """
    issues: list[dict] = []
    try:
        for slide in (getattr(spec, "slides", []) or []):
            layout = getattr(slide, "layout", "") or ""
            if layout == "section_header":
                continue
            notes = (getattr(slide, "speaker_notes", "") or "").strip()
            if notes and len(notes) < 60:
                issues.append({
                    "severity": "p2",
                    "check": "pptx_speaker_notes_too_short",
                    "message": (
                        f"幻灯片「{slide.assertion_title[:30]}」演讲者备注过短"
                        f"（{len(notes)} 字，建议至少60字），"
                        "请补充背景说明和衔接句"
                    ),
                    "section_id": getattr(slide, "id", None),
                })
    except Exception as exc:
        logger.debug("[QA] pptx_speaker_notes_length check error (non-fatal): %s", exc)
    return issues


def _check_pptx_agenda_sync(spec) -> list[dict]:
    """Check 23: R16-P2 PPTX agenda slide bullets should match actual slide titles.

    When an agenda/TOC slide is auto-inserted (or provided by the user), its bullet
    list should reflect the actual deck content.  If more than half the agenda items
    don't match any slide title, the agenda is misleading.

    Matching uses a prefix-substring heuristic (first 12 chars of each bullet must
    appear in some slide title, or vice versa) to tolerate minor wording differences.
    """
    _AGENDA_KEYWORDS = {"议程", "目录", "agenda", "outline", "contents", "overview", "概览", "今日"}
    issues: list[dict] = []
    try:
        slides = list(getattr(spec, "slides", []) or [])
        if len(slides) < 4:
            return issues

        # Find an agenda slide within the first 3 slides
        agenda_slide = None
        for s in slides[:3]:
            title_lower = (getattr(s, "assertion_title", "") or "").lower()
            if any(kw in title_lower for kw in _AGENDA_KEYWORDS):
                agenda_slide = s
                break

        if not agenda_slide:
            return issues

        agenda_bullets = [
            b.strip() for b in (getattr(agenda_slide, "bullets", []) or [])
            if (b or "").strip()
        ]
        if not agenda_bullets:
            return issues  # empty agenda doesn't make sense but can't compare

        # Build set of all other slide titles for lookup
        other_titles = [
            (getattr(s, "assertion_title", "") or "")
            for s in slides
            if s is not agenda_slide
        ]

        mismatched: list[str] = []
        for bullet in agenda_bullets:
            prefix = bullet[:12]
            # Check if bullet prefix appears in any slide title, or a slide title
            # prefix appears in the bullet (bi-directional substring match)
            found = any(
                prefix in title or title[:12] in bullet
                for title in other_titles
            )
            if not found:
                mismatched.append(bullet[:30])

        # Only flag when majority of agenda items are unmatched (>50%) — minor drift is OK
        if len(mismatched) > len(agenda_bullets) / 2:
            sample = mismatched[0] if mismatched else ""
            issues.append({
                "severity": "p2",
                "check": "pptx_agenda_mismatch",
                "message": (
                    f"议程页条目与幻灯片标题不匹配（{len(mismatched)}/{len(agenda_bullets)} 项），"
                    f"例：「{sample}」未在后续幻灯片标题中找到对应内容，"
                    "请同步议程与实际内容"
                ),
                "section_id": getattr(agenda_slide, "id", None),
            })
    except Exception as exc:
        logger.debug("[QA] pptx_agenda_sync check error (non-fatal): %s", exc)
    return issues


def _check_xlsx_formula_cross_refs(rendered_bytes: bytes) -> list[dict]:
    """Check 24: R17-P1 Validate cross-sheet formula references in rendered XLSX.

    Scans all cells for formulas containing ='SheetName'!ref references and verifies
    that each referenced SheetName actually exists as a worksheet tab in the workbook.
    Missing sheet references produce silent #REF! errors when the user opens the file.

    Only quoted sheet-name refs (='Name'!) are checked — bare refs are rarely
    generated by the spec_gen prompts and are harder to reliably parse.
    """
    issues: list[dict] = []
    try:
        import io as _io
        from openpyxl import load_workbook

        wb = load_workbook(_io.BytesIO(rendered_bytes), read_only=True, data_only=False)
        sheet_names: set[str] = {ws.title for ws in wb.worksheets}

        # ='SheetName'!ref  (quoted — handles names with spaces/CJK)
        _QUOTED_SHEET_REF_PAT = re.compile(r"='([^']+)'!", re.IGNORECASE)

        bad_refs: set[str] = set()
        for ws in wb.worksheets:
            for row in ws.iter_rows():
                for cell in row:
                    val = cell.value
                    if not isinstance(val, str) or not val.startswith("="):
                        continue
                    for m in _QUOTED_SHEET_REF_PAT.finditer(val):
                        ref_sheet = m.group(1)
                        if ref_sheet not in sheet_names:
                            bad_refs.add(ref_sheet)

        if bad_refs:
            sample = "、".join(sorted(bad_refs)[:3])
            issues.append({
                "severity": "p1",
                "check": "xlsx_formula_cross_ref",
                "message": (
                    f"Excel 公式中引用了不存在的工作表（{sample}），"
                    "这些引用将产生 #REF! 错误，请检查跨表公式中的工作表名称"
                ),
                "section_id": None,
            })
    except Exception as exc:
        logger.debug("[QA] xlsx_formula_cross_ref check error (non-fatal): %s", exc)
    return issues


def _check_docx_section_char_dropoff(spec, text: str) -> list[dict]:
    """Check 25: R17-P2 Detect LLM energy drop-off in the final sections of a DOCX.

    LLMs sometimes generate increasingly thin content toward the end of a long document.
    This manifests as the last 3 sections all falling below 40% of their target_chars.
    Two or more such sections in a row indicates systematic under-generation, not
    just one short section — that would be caught by per_section_chars (Check 9).
    """
    issues: list[dict] = []
    try:
        sections = list(getattr(spec, "sections", []) or [])
        if len(sections) < 5:
            return issues  # too short to have a meaningful "tail" pattern

        last_sections = sections[-3:]
        title_list = [(s.title, s.id, getattr(s, "target_chars", 400)) for s in sections]

        short_count = 0
        for sec in last_sections:
            title_pos = text.find(sec.title)
            if title_pos < 0:
                continue
            target = getattr(sec, "target_chars", 400)
            # Find start of next section to bound this section's body text
            next_pos = len(text)
            for other_title, _, _ in title_list:
                pos = text.find(other_title, title_pos + len(sec.title))
                if 0 < pos < next_pos:
                    next_pos = pos
            body = text[title_pos + len(sec.title):next_pos]
            actual = len(re.sub(r"\s+", "", body))
            if actual < target * 0.40:
                short_count += 1

        if short_count >= 2:
            issues.append({
                "severity": "p2",
                "check": "docx_section_char_dropoff",
                "message": (
                    f"文档末尾 {short_count}/3 个章节字数明显低于目标（<40%），"
                    "可能存在 LLM 生成能量衰减，建议重新生成尾部章节"
                ),
                "section_id": None,
            })
    except Exception as exc:
        logger.debug("[QA] docx_section_char_dropoff check error (non-fatal): %s", exc)
    return issues


def _check_docx_evidence_coverage(spec, intent: str, has_research: bool) -> list[dict]:
    """Check 26: R17-P2 fill_from_reference DOCX sections should cite evidence sources.

    When the pipeline ran research and the intent is reference-based, each section
    should have at least one entry in evidence_ids.  Sections with empty evidence_ids
    likely contain hallucinated content not grounded in retrieved documents.

    Only fires when the majority (>50%) of sections are uncited, to avoid flagging
    documents where the research pipeline returned nothing for some sections.
    """
    issues: list[dict] = []
    if intent not in ("fill_from_reference", "extend"):
        return issues
    if not has_research:
        return issues  # no research was run — no expectation of citation
    try:
        sections = list(getattr(spec, "sections", []) or [])
        if not sections:
            return issues

        uncited = [s for s in sections if not getattr(s, "evidence_ids", None)]
        if len(uncited) > len(sections) / 2:
            sample = uncited[0].title[:20] if uncited else ""
            issues.append({
                "severity": "p2",
                "check": "docx_evidence_coverage",
                "message": (
                    f"{len(uncited)}/{len(sections)} 个章节 evidence_ids 为空"
                    f"（例：「{sample}」），fill_from_reference 文档应标注数据来源，"
                    "请确认研究材料是否被正确关联"
                ),
                "section_id": None,
            })
    except Exception as exc:
        logger.debug("[QA] docx_evidence_coverage check error (non-fatal): %s", exc)
    return issues


def _check_docx_last_section_closing(spec, text: str) -> list[dict]:
    """Check 27: R18-P2 The final DOCX section should end with a closing/summary sentence.

    A well-written document closes with a conclusion sentence that wraps up the content.
    This check looks for closing-sentence keywords in the last paragraph of the final
    non-empty section.  Complements Check 21 (transition sentences for non-last sections).

    Only fires for documents with ≥3 sections, since short docs may not need formal closing.
    """
    _CLOSING_KEYWORDS = (
        "总结", "综上", "总体", "小结", "结语", "结论", "结束", "综合", "最终",
        "总而言之", "综合来看", "回顾全文", "总的来说",
        "展望", "下一步", "期待", "未来",
    )
    issues: list[dict] = []
    try:
        sections = [s for s in (getattr(spec, "sections", []) or []) if s.paragraphs]
        if len(sections) < 3:
            return issues  # too short to require a formal closing

        last_sec = sections[-1]
        last_para = (last_sec.paragraphs[-1] if last_sec.paragraphs else "").strip()
        if not last_para:
            return issues

        has_closing = any(kw in last_para for kw in _CLOSING_KEYWORDS)
        if not has_closing:
            issues.append({
                "severity": "p2",
                "check": "docx_last_section_closing",
                "message": (
                    f"末章节「{last_sec.title[:25]}」的最后一段缺少收尾/总结句，"
                    "文档应以「综上所述」「总体来看」等收官语结束，"
                    "避免在分析中途突然截断"
                ),
                "section_id": last_sec.id,
            })
    except Exception as exc:
        logger.debug("[QA] docx_last_section_closing check error (non-fatal): %s", exc)
    return issues


def _check_xlsx_named_range_conflicts(rendered_bytes: bytes) -> list[dict]:
    """Check 28: R18-P1 Detect named range conflicts in rendered Excel workbook.

    A named range conflict occurs when:
    1. A workbook-scoped name shadows a sheet-scoped name with the same text, OR
    2. Two sheets define the same local named range (causing confusion when cross-referenced)

    These conflicts are invisible until the user tries to use the workbook's name manager
    or references a name across sheets — catching them at QA time prevents silent errors.
    """
    issues: list[dict] = []
    try:
        import io
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(rendered_bytes))

        # Collect all defined names (workbook-level)
        defined_names = list(wb.defined_names.definedName) if hasattr(wb.defined_names, "definedName") else []

        # openpyxl ≥3.1 uses .keys() on defined_names; fallback to iteration
        try:
            name_keys = list(wb.defined_names.keys())
        except Exception:
            name_keys = [getattr(dn, "name", "") for dn in defined_names]

        # Group by name (case-insensitive)
        from collections import Counter
        name_counts = Counter(n.upper() for n in name_keys if n)
        duplicates = [n for n, count in name_counts.items() if count > 1]

        if duplicates:
            sample = duplicates[0]
            issues.append({
                "severity": "p1",
                "check": "xlsx_named_range_conflict",
                "message": (
                    f"Excel 工作簿中存在重复命名的定义名称（{len(duplicates)} 个），"
                    f"例：「{sample}」，重名会导致 Name Manager 冲突和公式引用歧义，"
                    "请确保每个命名范围在工作簿中唯一"
                ),
                "section_id": None,
            })

        # Also check for worksheet-scoped conflicts with workbook-scoped names
        sheet_names_upper = {ws.title.upper() for ws in wb.worksheets}
        name_shadows_sheet = [
            n for n in name_keys
            if n and n.upper() in sheet_names_upper and len(n) >= 2
        ]
        if name_shadows_sheet:
            sample = name_shadows_sheet[0]
            issues.append({
                "severity": "p2",
                "check": "xlsx_named_range_shadows_sheet",
                "message": (
                    f"命名范围「{sample}」与工作表名称相同，"
                    "可能导致 ='SheetName'! 跨表引用产生歧义，"
                    "建议为命名范围使用不同于工作表的名称"
                ),
                "section_id": None,
            })
    except Exception as exc:
        logger.debug("[QA] xlsx_named_range_conflicts check error (non-fatal): %s", exc)
    return issues
