import logging
import os
from pathlib import Path

from app.config import settings

logger = logging.getLogger(__name__)


# ── Encoding detection (SOTA: handles Chinese GBK/GB18030/Big5 + UTF-8 BOM) ──
# This addresses the cross-cutting weakness that broke all Chinese plaintext
# inputs on Windows-exported files. We try chardet first, then a curated
# fallback chain that covers >99% of real-world Chinese / Japanese / Korean files.

_ENCODING_FALLBACK_CHAIN = (
    "utf-8-sig",   # UTF-8 with BOM (common Windows export)
    "utf-8",
    "gb18030",     # superset of GBK/GB2312, covers all Chinese
    "big5",        # Traditional Chinese
    "shift_jis",   # Japanese
    "euc-kr",      # Korean
    "latin-1",     # last-resort: never throws
)


def detect_encoding(raw_bytes: bytes, sample_size: int = 65536) -> str:
    """Detect the most likely encoding for a byte stream.

    Strategy:
      1. Check for BOM markers (most reliable)
      2. Try chardet on first sample_size bytes (fast + accurate for CJK)
      3. Validate by attempting decode; on failure walk the fallback chain
    """
    if not raw_bytes:
        return "utf-8"

    # 1) BOM detection
    if raw_bytes.startswith(b"\xef\xbb\xbf"):
        return "utf-8-sig"
    if raw_bytes.startswith(b"\xff\xfe") or raw_bytes.startswith(b"\xfe\xff"):
        return "utf-16"

    sample = raw_bytes[:sample_size]

    # 2) chardet
    try:
        import chardet
        detected = chardet.detect(sample)
        enc = (detected.get("encoding") or "").lower()
        conf = float(detected.get("confidence") or 0)
        if enc and conf >= 0.6:
            # Normalize common aliases
            if enc in ("gb2312", "gbk"):
                enc = "gb18030"  # superset, always safer
            # Validate
            try:
                sample.decode(enc)
                return enc
            except (UnicodeDecodeError, LookupError):
                pass
    except ImportError:
        logger.debug("chardet not available, using fallback chain")
    except Exception as exc:
        logger.debug("chardet failed: %s", exc)

    # 3) Fallback chain
    for enc in _ENCODING_FALLBACK_CHAIN:
        try:
            sample.decode(enc)
            return enc
        except (UnicodeDecodeError, LookupError):
            continue

    return "latin-1"  # never throws


def read_text_file_smart(path: Path | str) -> str:
    """Read a text file with automatic encoding detection.

    Replaces the dangerous `open(..., 'utf-8', errors='replace')` pattern
    that silently corrupted Chinese GBK files into question marks.
    """
    p = Path(path)
    with open(p, "rb") as f:
        raw = f.read()
    enc = detect_encoding(raw)
    try:
        return raw.decode(enc)
    except Exception:
        # Last resort: lossy decode
        return raw.decode(enc, errors="replace")


async def extract_text(file_path: str, file_type: str) -> str:
    """Extract text content from uploaded files."""
    resolved_path = _resolve_file_path(file_path)
    ext = resolved_path.suffix.lower()
    text = ""

    try:
        if ext == ".pdf":
            text = _extract_pdf(str(resolved_path))
        elif ext in (".docx", ".doc"):
            text = _extract_docx(str(resolved_path))
        elif ext in (".xlsx", ".xls"):
            text = _extract_xlsx(str(resolved_path))
        elif ext in (".pptx", ".ppt"):
            text = _extract_pptx(str(resolved_path))
        elif ext in (".csv", ".tsv"):
            text = _extract_delimited_table(str(resolved_path), delimiter="\t" if ext == ".tsv" else ",")
        elif ext in (".txt", ".md", ".json", ".log", ".yaml", ".yml", ".xml", ".html"):
            text = read_text_file_smart(resolved_path)
        else:
            text = f"[Binary file: {ext}]"
    except Exception as e:
        logger.warning("File extraction failed for %s: %s", resolved_path, e)
        text = f"[Extraction failed: {e}]"

    return text[:100000]  # Truncate to 100k chars


def _resolve_file_path(file_path: str) -> Path:
    """Resolve legacy relative upload paths across old and new data dirs."""
    path = Path(file_path)
    if path.exists():
        return path
    if not path.is_absolute():
        candidates = [
            Path.cwd() / path,
            Path(settings.upload_dir) / path.name,
            Path("data/uploads") / path.name,
            Path("uploads") / path.name,
        ]
        for candidate in candidates:
            if candidate.exists():
                return candidate
    return path


def _extract_pdf(path: str) -> str:
    parts: list[str] = []
    table_count = 0
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            meta = pdf.metadata or {}
            title = meta.get("Title", "") or meta.get("Subject", "")
            if title:
                parts.append(f"[PDF_TITLE: {title}]")
            for i, page in enumerate(pdf.pages[:50]):
                page_parts: list[str] = []
                # Text
                t = page.extract_text()
                if t:
                    page_parts.append(t)
                # Tables — convert to markdown
                try:
                    tables = page.extract_tables() or []
                    for table in tables:
                        if not table:
                            continue
                        table_count += 1
                        md: list[str] = []
                        for ri, row in enumerate(table[:60]):
                            cells = [str(c or "").strip().replace("\n", " ") for c in row]
                            md.append("| " + " | ".join(cells) + " |")
                            if ri == 0:
                                md.append("|" + "|".join(["---"] * len(cells)) + "|")
                        page_parts.append(f"\n[PDF_TABLE {table_count}]\n" + "\n".join(md))
                except Exception:
                    pass
                if page_parts:
                    parts.append(f"--- Page {i + 1} ---")
                    parts.extend(page_parts)
    except Exception:
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(path)
            for page in reader.pages[:50]:
                t = page.extract_text()
                if t:
                    parts.append(t)
        except Exception:
            return "[PDF parse error]"
    if table_count:
        parts.append(f"[PDF_TABLES: count={table_count}]")

    # Enhanced image extraction via PyMuPDF
    image_parts = _extract_pdf_images_fitz(path)
    parts.extend(image_parts)

    return "\n".join(parts)


def _extract_pdf_images_fitz(path: str) -> list[str]:
    """Extract images and detect vector charts in a PDF using PyMuPDF.

    For each page:
    - Embedded raster images: extract dimensions, try OCR via pytesseract.
    - Pages with complex vector drawings (charts): record drawing count and bounding box.
    """
    results: list[str] = []
    total_images = 0

    try:
        import fitz  # PyMuPDF
        from PIL import Image
        import io

        doc = fitz.open(path)

        for page_num in range(min(len(doc), 50)):
            page = doc[page_num]

            # ── Embedded raster images ────────────────────────────────────────
            img_list = page.get_images(full=True)
            for img_info in img_list[:8]:  # cap at 8 per page
                try:
                    xref = img_info[0]
                    base_img = doc.extract_image(xref)
                    img_bytes = base_img["image"]
                    img_w = base_img.get("width", 0)
                    img_h = base_img.get("height", 0)
                    img_ext = base_img.get("ext", "png")

                    # Skip tiny decorative icons
                    if img_w < 80 or img_h < 80:
                        continue

                    total_images += 1
                    ocr_text = ""
                    type_hint = ""

                    try:
                        pil_img = Image.open(io.BytesIO(img_bytes)).convert("RGB")

                        # Infer basic image type from color diversity
                        type_hint = _infer_pdf_image_type(pil_img)

                        # OCR — gracefully skip if tesseract not installed
                        try:
                            import pytesseract
                            raw = pytesseract.image_to_string(
                                pil_img, lang="chi_sim+eng", config="--psm 6"
                            ).strip()
                            ocr_text = " ".join(raw.split())[:400]
                        except Exception:
                            pass
                    except Exception:
                        pass

                    marker = (
                        f"[PDF_IMAGE: page={page_num + 1}"
                        f" size={img_w}x{img_h} ext={img_ext}"
                    )
                    if type_hint:
                        marker += f" | type={type_hint}"
                    if ocr_text:
                        marker += f" | ocr_text={ocr_text}"
                    marker += "]"
                    results.append(marker)
                except Exception:
                    continue

            # ── Vector-drawn charts/diagrams ──────────────────────────────────
            # Pages without embedded images but with many colored vector paths
            # are almost certainly rendered charts (e.g. bar/pie/line charts).
            if not img_list:
                try:
                    drawings = page.get_drawings()
                    colored = [
                        d for d in drawings
                        if d.get("fill") and d["fill"] not in (None, (1, 1, 1), (1.0, 1.0, 1.0))
                    ]
                    if len(colored) >= 12:
                        # Derive bounding box of all drawing paths
                        rects = [d["rect"] for d in colored if d.get("rect")]
                        if rects:
                            x0 = min(r.x0 for r in rects)
                            y0 = min(r.y0 for r in rects)
                            x1 = max(r.x1 for r in rects)
                            y1 = max(r.y1 for r in rects)
                            bbox = f"{int(x1 - x0)}x{int(y1 - y0)}pt"
                        else:
                            bbox = "unknown"
                        results.append(
                            f"[PDF_VECTOR_CHART: page={page_num + 1}"
                            f" colored_paths={len(colored)} bbox={bbox}"
                            f" — 疑似矢量图表/示意图]"
                        )
                except Exception:
                    pass

        doc.close()

    except Exception:
        pass

    if total_images:
        results.append(f"[PDF_IMAGES: count={total_images}]")

    return results


def _infer_pdf_image_type(pil_img) -> str:
    """Heuristic image-type inference based on color statistics (no heavy deps)."""
    try:
        # Sample a small thumbnail for speed
        thumb = pil_img.resize((64, 64))
        pixels = list(thumb.getdata())
        unique = len(set(pixels))
        total = len(pixels)

        if unique > total * 0.5:
            return "照片/复杂图像"
        if unique < 20:
            return "纯色图形/图标"

        # Check for strong horizontal banding (bar charts, timelines)
        row_brightness = []
        w, h = thumb.size
        for y in range(h):
            row_px = [thumb.getpixel((x, y)) for x in range(w)]
            avg = sum(sum(p[:3]) / 3 for p in row_px) / w
            row_brightness.append(avg)
        variance = sum((v - sum(row_brightness) / h) ** 2 for v in row_brightness) / h

        if variance > 800:
            return "柱状图/条形图"
        if unique < 80:
            return "图表/示意图"
        return "图示"
    except Exception:
        return ""


def _extract_docx(path: str) -> str:
    try:
        from docx import Document
        doc = Document(path)
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table_idx, table in enumerate(doc.tables, start=1):
            rows = []
            for row in table.rows:
                cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                if any(cells):
                    rows.append(" | ".join(cells))
            if rows:
                parts.append(f"--- Table {table_idx} ---")
                parts.extend(rows)
        # Detect embedded images and charts via inline_shapes
        image_count = 0
        chart_count = 0
        try:
            for shape in doc.inline_shapes:
                # WD_INLINE_SHAPE: PICTURE=3, CHART=5
                if shape.type == 3:
                    image_count += 1
                elif shape.type == 5:
                    chart_count += 1
                elif shape.type not in (3, 4, 5, 6):
                    # Unknown types — count as images conservatively
                    image_count += 1
        except Exception:
            pass
        # Fallback: scan XML for chart and drawing references
        try:
            import re as _re
            xml_str = doc.element.xml
            xml_chart_count = len(_re.findall(r'<c:chart\b|<c:chartSpace\b', xml_str))
            xml_image_count = len(_re.findall(r'<a:blip\b', xml_str))
            if xml_chart_count > chart_count:
                chart_count = xml_chart_count
            if xml_image_count > image_count:
                image_count = xml_image_count
        except Exception:
            pass
        if chart_count:
            parts.append(f"[EMBEDDED_CHARTS: count={chart_count}]")
        if image_count:
            parts.append(f"[EMBEDDED_IMAGES: count={image_count}]")
        return "\n".join(parts)
    except Exception:
        return "[DOCX parse error]"


_XLSX_CHART_TYPE_MAP = {
    "BarChart": "柱状对比图", "LineChart": "折线趋势图", "PieChart": "饼图/分布图",
    "AreaChart": "面积图", "ScatterChart": "散点图", "BubbleChart": "气泡图",
    "RadarChart": "雷达图", "DoughnutChart": "环形图", "StockChart": "股价图",
    "SurfaceChart": "曲面图",
}


def _extract_xlsx(path: str) -> str:
    """Extract XLSX with full numeric profiling and multi-sheet support."""
    parts: list[str] = []
    total_charts = 0

    # ── Phase 1: pandas multi-sheet read + profile ──────────────────────────
    try:
        import pandas as pd
        xl = pd.ExcelFile(path, engine="openpyxl")
        sheet_names = xl.sheet_names
        parts.append(f"[EXCEL_SHEETS: {len(sheet_names)} 个工作表: {', '.join(sheet_names[:12])}]")

        for sheet_name in sheet_names[:6]:            # profile up to 6 sheets
            try:
                df = xl.parse(sheet_name, header=0)
                n_rows, n_cols = df.shape
                parts.append(f"--- Sheet: {sheet_name} ({n_rows}行 × {n_cols}列) ---")

                # Column schema
                dtypes_str = ", ".join(f"{c}({t})" for c, t in df.dtypes.items())
                parts.append(f"[列结构] {dtypes_str}")

                # Null counts for columns with nulls
                null_counts = df.isnull().sum()
                nulls_nonzero = {str(k): int(v) for k, v in null_counts.items() if v > 0}
                if nulls_nonzero:
                    parts.append(f"[空值] {nulls_nonzero}")

                # Numeric profile
                numeric_df = df.select_dtypes(include="number")
                if not numeric_df.empty:
                    desc = numeric_df.describe().round(4)
                    parts.append("[数值列统计]")
                    parts.append(desc.to_string())

                    # Sum & total for common finance fields
                    sums = numeric_df.sum(numeric_only=True).round(4)
                    parts.append("[列合计] " + " | ".join(
                        f"{col}={v:,.4f}" for col, v in sums.items()
                    ))

                # Sample rows (up to 50 for small, 20 for large sheets)
                sample_n = 20 if n_rows > 200 else min(n_rows, 50)
                parts.append(f"[前{sample_n}行数据]")
                parts.append(df.head(sample_n).to_string(index=False))

                # If sheet is large, also show last 5 rows (totals often there)
                if n_rows > sample_n + 5:
                    parts.append("[末5行数据（可能含合计行）]")
                    parts.append(df.tail(5).to_string(index=False))

                    # Show ALL data for medium sheets (≤1000 rows)
                    if n_rows <= 1000:
                        parts.append(f"[全部{n_rows}行数据]")
                        parts.append(df.to_string(index=False))

            except Exception as sheet_err:
                parts.append(f"--- Sheet: {sheet_name} [读取失败: {sheet_err}] ---")

    except Exception as pandas_err:
        logger.debug("pandas ExcelFile failed for %s: %s", path, pandas_err)

    # ── Phase 2: openpyxl for chart metadata + overflow rows ───────────────
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, data_only=True)

        for name in wb.sheetnames[:6]:
            ws = wb[name]

            # If pandas already handled it, just add charts
            charts = getattr(ws, "_charts", [])
            for chart in charts:
                total_charts += 1
                type_cn = _XLSX_CHART_TYPE_MAP.get(type(chart).__name__, type(chart).__name__)
                title_text = ""
                try:
                    if chart.title:
                        title_text = str(chart.title)
                except Exception:
                    pass
                series_names: list[str] = []
                try:
                    series_names = [
                        str(s.title) for s in chart.series if getattr(s, "title", None)
                    ][:4]
                except Exception:
                    pass
                marker = f"[CHART: type={type_cn}"
                if title_text:
                    marker += f" | title={title_text}"
                if series_names:
                    marker += f" | series={', '.join(series_names)}"
                marker += "]"
                parts.append(marker)

            # Named ranges
            try:
                named_ranges = list(wb.defined_names.keys())[:10]
                if named_ranges:
                    parts.append(f"[命名区域] {', '.join(named_ranges)}")
            except Exception:
                pass

    except Exception as openpyxl_err:
        logger.debug("openpyxl fallback failed for %s: %s", path, openpyxl_err)
        # Ultimate fallback: just return whatever we have
        if not parts:
            return "[XLSX parse error]"

    if total_charts:
        parts.append(f"[XLSX_TOTAL_CHARTS: {total_charts}]")

    return "\n".join(parts) if parts else "[XLSX parse error]"


def _extract_delimited_table(path: str, delimiter: str = ",") -> str:
    """Extract CSV/TSV with an advanced columnar summary when available."""
    # Smart encoding detection — critical for Chinese GBK CSVs exported from
    # Excel for Windows (most common real-world failure mode)
    with open(path, "rb") as f:
        raw = f.read()
    encoding = detect_encoding(raw)

    fallback = ""
    try:
        fallback = raw.decode(encoding, errors="replace")[:60000]
    except Exception:
        try:
            fallback = raw.decode("utf-8", errors="replace")[:60000]
        except Exception:
            fallback = ""

    try:
        import polars as pl

        df = pl.read_csv(
            path,
            separator=delimiter,
            infer_schema_length=1000,
            ignore_errors=True,
            encoding=encoding if encoding in ("utf-8", "utf-8-lossy") else "utf8-lossy",
        )
        n_rows, n_cols = df.height, df.width
        parts = [
            "--- Columnar Data Summary (Polars) ---",
            f"rows={n_rows}, columns={n_cols}",
            "schema=" + ", ".join(f"{name}:{dtype}" for name, dtype in df.schema.items()),
        ]
        nulls = df.null_count().to_dicts()[0] if n_cols else {}
        if nulls:
            parts.append("null_count=" + ", ".join(f"{k}:{v}" for k, v in nulls.items() if v))

        numeric_cols = [
            name for name, dtype in df.schema.items()
            if dtype.is_numeric()
        ]
        if numeric_cols:
            parts.append("--- Numeric Profile (mean | sum | min | max | std) ---")
            exprs = []
            for name in numeric_cols[:16]:
                exprs.extend([
                    pl.col(name).mean().round(4).alias(f"{name}_mean"),
                    pl.col(name).sum().round(4).alias(f"{name}_sum"),
                    pl.col(name).min().alias(f"{name}_min"),
                    pl.col(name).max().alias(f"{name}_max"),
                    pl.col(name).std().round(4).alias(f"{name}_std"),
                ])
            profile = df.select(exprs)
            parts.append(str(profile))

        # Sample rows — more rows for smaller files
        sample_n = 20 if n_rows > 500 else min(n_rows, 100)
        parts.append(f"--- Sample Rows (first {sample_n}) ---")
        parts.append(str(df.head(sample_n)))

        # For medium files, include all data
        if n_rows <= 500:
            parts.append(f"--- All {n_rows} Rows ---")
            parts.append(str(df))
        elif n_rows <= 2000:
            parts.append("--- Last 10 Rows ---")
            parts.append(str(df.tail(10)))

        if fallback and n_rows > 2000:
            parts.append("--- Raw Preview (large file) ---")
            parts.append(fallback[:15000])
        return "\n".join(parts)
    except Exception:
        try:
            import duckdb

            rel = duckdb.read_csv(path, delim=delimiter)
            desc = rel.describe().fetchdf()
            sample = rel.limit(30).fetchdf()
            return "\n".join([
                "--- Columnar Data Summary (DuckDB) ---",
                str(desc),
                "--- Sample Rows ---",
                str(sample),
                "--- Raw Preview ---",
                fallback[:20000],
            ])
        except Exception:
            return fallback or "[Delimited table parse error]"


def _extract_pptx(path: str) -> str:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        _CHART_TYPE_MAP = {
            "BAR_CLUSTERED": "柱状对比图", "BAR_STACKED": "柱状堆叠图",
            "BAR_STACKED_100": "百分比柱状图", "COLUMN_CLUSTERED": "柱状对比图",
            "COLUMN_STACKED": "柱状堆叠图", "LINE": "折线趋势图",
            "LINE_MARKERS": "折线趋势图", "PIE": "饼图/分布图",
            "PIE_EXPLODED": "饼图/分布图", "DOUGHNUT": "环形图",
            "AREA": "面积图", "SCATTER": "散点图", "BUBBLE": "气泡图",
            "RADAR": "雷达图", "XY_SCATTER": "散点图",
        }

        prs = Presentation(path)
        parts = []
        image_count = 0
        chart_count = 0
        for i, slide in enumerate(prs.slides):
            parts.append(f"--- Slide {i+1} ---")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        parts.append(text)
                if shape.has_chart:
                    chart = shape.chart
                    chart_count += 1
                    type_key = str(chart.chart_type).split(".")[-1] if chart.chart_type else ""
                    type_cn = _CHART_TYPE_MAP.get(type_key, type_key or "图表")
                    chart_title = ""
                    try:
                        if chart.has_title:
                            chart_title = chart.chart_title.text_frame.text.strip()
                    except Exception:
                        pass
                    # Extract series names AND actual data values
                    series_lines: list[str] = []
                    categories: list[str] = []
                    try:
                        for s in chart.series[:6]:
                            sname = s.name or "系列"
                            try:
                                vals = [round(v, 4) for v in s.values if v is not None][:12]
                            except Exception:
                                vals = []
                            if vals:
                                series_lines.append(f"{sname}={vals}")
                            else:
                                series_lines.append(sname)
                    except Exception:
                        pass
                    # Try to extract category labels from chart XML
                    try:
                        from lxml import etree as _et
                        nsmap = {"c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
                                 "a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
                        chart_elm = chart._element
                        cat_nodes = chart_elm.findall(".//c:cat//c:v", nsmap) or \
                                    chart_elm.findall(".//c:xVal//c:v", nsmap)
                        categories = [n.text for n in cat_nodes if n.text][:12]
                    except Exception:
                        pass
                    marker = f"[CHART: type={type_cn}"
                    if chart_title:
                        marker += f" | title={chart_title}"
                    if categories:
                        marker += f" | categories={categories}"
                    if series_lines:
                        marker += f" | data={'; '.join(series_lines)}"
                    marker += "]"
                    parts.append(marker)
                elif hasattr(shape, "shape_type") and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_count += 1
                    parts.append(f"[IMAGE: {shape.name}]")
            # Speaker notes
            try:
                notes_frame = slide.notes_slide.notes_text_frame
                notes_text = notes_frame.text.strip() if notes_frame else ""
                if notes_text:
                    parts.append(f"[NOTES: {notes_text[:300]}]")
            except Exception:
                pass
        if chart_count:
            parts.append(f"[PPTX_TOTAL_CHARTS: {chart_count}]")
        if image_count:
            parts.append(f"[PPTX_TOTAL_IMAGES: {image_count}]")
        return "\n".join(parts)
    except Exception:
        return "[PPTX parse error]"
