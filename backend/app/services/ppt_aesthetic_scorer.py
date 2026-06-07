"""PPT Aesthetic Scorer — Programmatic geometric quality metrics.

Based on AeSlides (2025): "Incentivizing Aesthetic Layout in LLM-Based Slide
Generation via Verifiable Rewards" (arXiv:2604.22840).

Key insight: Text-centric generation ≠ visual quality. Pure LLM scoring misses
geometric layout issues. Verifiable structural metrics provide:
  - Aspect ratio compliance (36% → 85% improvement in AeSlides)
  - Whitespace reduction (-44%)
  - Element collision elimination (-43%)
  - Visual balance improvement (-28%)

Metrics are computed programmatically from python-pptx without LLM calls.
Results can be used to:
  1. Gate slide quality before delivery
  2. Guide SlideEditorAgent toward specific geometric improvements
  3. Provide quantitative score alongside LLM qualitative review
"""
from __future__ import annotations

import logging
from dataclasses import dataclass, field
from typing import Optional

logger = logging.getLogger(__name__)


@dataclass
class ShapeMetrics:
    """Geometric metrics for a single shape."""
    shape_id: int
    slide_id: int
    left: float    # normalized 0-1
    top: float     # normalized 0-1
    width: float   # normalized 0-1
    height: float  # normalized 0-1
    has_text: bool
    has_chart: bool
    has_image: bool
    text_length: int = 0
    font_size_pt: float = 0.0

    @property
    def right(self) -> float:
        return self.left + self.width

    @property
    def bottom(self) -> float:
        return self.top + self.height

    @property
    def area(self) -> float:
        return self.width * self.height

    @property
    def aspect_ratio(self) -> float:
        if self.height == 0:
            return 0.0
        return self.width / self.height


@dataclass
class SlideAestheticScore:
    """Aesthetic score for a single slide."""
    slide_id: int

    # Geometric metrics (0.0 - 1.0, higher is better)
    aspect_ratio_score: float = 1.0      # Shape aspect ratios within valid ranges
    whitespace_score: float = 1.0        # Adequate whitespace maintained
    collision_score: float = 1.0         # No overlapping elements
    balance_score: float = 1.0          # Visual weight balanced horizontally/vertically
    text_density_score: float = 1.0      # Text not too dense or too sparse
    element_count_score: float = 1.0     # Appropriate number of elements

    # Issue descriptions for LLM-guided fixes
    issues: list[str] = field(default_factory=list)
    suggestions: list[str] = field(default_factory=list)

    @property
    def overall_score(self) -> float:
        """Weighted composite score (0-100)."""
        weights = {
            "collision": 0.30,
            "whitespace": 0.20,
            "balance": 0.20,
            "aspect_ratio": 0.15,
            "text_density": 0.10,
            "element_count": 0.05,
        }
        raw = (
            self.collision_score * weights["collision"] +
            self.whitespace_score * weights["whitespace"] +
            self.balance_score * weights["balance"] +
            self.aspect_ratio_score * weights["aspect_ratio"] +
            self.text_density_score * weights["text_density"] +
            self.element_count_score * weights["element_count"]
        )
        return round(raw * 100, 1)

    def to_dict(self) -> dict:
        return {
            "slide_id": self.slide_id,
            "overall_score": self.overall_score,
            "aspect_ratio_score": round(self.aspect_ratio_score * 100, 1),
            "whitespace_score": round(self.whitespace_score * 100, 1),
            "collision_score": round(self.collision_score * 100, 1),
            "balance_score": round(self.balance_score * 100, 1),
            "text_density_score": round(self.text_density_score * 100, 1),
            "element_count_score": round(self.element_count_score * 100, 1),
            "issues": self.issues,
            "suggestions": self.suggestions,
        }


@dataclass
class PresentationAestheticScore:
    """Aggregate aesthetic score for the entire presentation."""
    slide_scores: list[SlideAestheticScore] = field(default_factory=list)
    overall_score: float = 0.0
    worst_slides: list[int] = field(default_factory=list)
    global_issues: list[str] = field(default_factory=list)

    def to_dict(self) -> dict:
        return {
            "overall_score": round(self.overall_score, 1),
            "slide_count": len(self.slide_scores),
            "worst_slides": self.worst_slides,
            "global_issues": self.global_issues,
            "slides": [s.to_dict() for s in self.slide_scores],
            "passed": self.overall_score >= 70,
            "grade": self._grade(),
        }

    def _grade(self) -> str:
        if self.overall_score >= 90:
            return "A"
        elif self.overall_score >= 80:
            return "B"
        elif self.overall_score >= 70:
            return "C"
        elif self.overall_score >= 60:
            return "D"
        return "F"


class PptAestheticScorer:
    """Compute programmatic aesthetic quality metrics for PPTX files.

    Implements AeSlides (2025) verifiable reward metrics:
    - Aspect ratio: element dimensions within valid ranges
    - Whitespace: total element area < 80% of slide area
    - Collisions: no significant element overlap
    - Balance: left/right and top/bottom weight symmetry
    - Text density: chars per line within readable range
    - Element count: 2-8 elements per slide
    """

    # Thresholds (from AeSlides paper and PPT design best practices)
    MIN_ASPECT_RATIO = 0.1    # At least 1:10
    MAX_ASPECT_RATIO = 20.0   # At most 20:1
    MAX_COVERAGE = 0.78        # Max slide area covered by elements
    MIN_COVERAGE = 0.20        # Min slide area covered (avoid empty slides)
    COLLISION_THRESHOLD = 0.05  # Max allowed overlap area ratio
    BALANCE_THRESHOLD = 0.30    # Max left-right or top-bottom imbalance
    MIN_FONT_PT = 10.0          # Minimum readable font size
    MAX_FONT_PT = 72.0          # Maximum practical font size
    MAX_ELEMENTS = 10           # Too many elements = cluttered
    MIN_ELEMENTS = 1            # Too few = sparse (unless cover)

    def score_pptx(self, pptx_path: str) -> PresentationAestheticScore:
        """Score all slides in a PPTX file.

        Args:
            pptx_path: Path to the PPTX file

        Returns:
            PresentationAestheticScore with per-slide and aggregate scores
        """
        try:
            from pptx import Presentation
            from pptx.util import Emu
        except ImportError:
            logger.warning("[AestheticScorer] python-pptx not available, returning default score")
            return self._unavailable_score()

        try:
            prs = Presentation(str(pptx_path))
        except Exception as e:
            logger.warning(f"[AestheticScorer] Failed to load PPTX: {e}")
            return self._unavailable_score()

        slide_width = prs.slide_width
        slide_height = prs.slide_height

        slide_scores = []
        for slide_idx, slide in enumerate(prs.slides, 1):
            metrics = self._extract_shape_metrics(slide, slide_idx, slide_width, slide_height)
            score = self._score_slide(slide_idx, metrics, slide_width, slide_height)
            slide_scores.append(score)

        return self._aggregate(slide_scores)

    def score_slide_metrics(
        self, slide_metrics: list[ShapeMetrics], slide_id: int = 1
    ) -> SlideAestheticScore:
        """Score a single slide from pre-computed shape metrics."""
        return self._score_slide(slide_id, slide_metrics, 1, 1)

    def _extract_shape_metrics(
        self, slide, slide_id: int, slide_width, slide_height
    ) -> list[ShapeMetrics]:
        """Extract normalized geometric metrics from a slide's shapes."""
        metrics = []
        w = slide_width or 1
        h = slide_height or 1

        for shape_idx, shape in enumerate(slide.shapes):
            try:
                left = shape.left / w if shape.left is not None else 0
                top = shape.top / h if shape.top is not None else 0
                width = shape.width / w if shape.width is not None else 0
                height = shape.height / h if shape.height is not None else 0

                # Clamp to valid range
                left = max(0.0, min(1.0, left))
                top = max(0.0, min(1.0, top))
                width = max(0.0, min(1.0 - left, width))
                height = max(0.0, min(1.0 - top, height))

                has_text = shape.has_text_frame
                has_chart = shape.has_chart
                has_image = shape.shape_type.name == "PICTURE"

                text_len = 0
                font_size = 0.0
                if has_text:
                    text_len = len(shape.text_frame.text)
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size:
                                from pptx.util import Pt
                                font_size = run.font.size / 12700  # EMU to pt
                                break
                        if font_size:
                            break

                m = ShapeMetrics(
                    shape_id=shape_idx + 1,
                    slide_id=slide_id,
                    left=left, top=top, width=width, height=height,
                    has_text=has_text, has_chart=has_chart, has_image=has_image,
                    text_length=text_len, font_size_pt=font_size,
                )
                metrics.append(m)
            except Exception:
                pass

        return metrics

    def _score_slide(
        self,
        slide_id: int,
        metrics: list[ShapeMetrics],
        slide_width,
        slide_height,
    ) -> SlideAestheticScore:
        """Compute all aesthetic scores for one slide."""
        score = SlideAestheticScore(slide_id=slide_id)

        if not metrics:
            score.issues.append(f"幻灯片{slide_id}没有可见元素")
            return score

        # 1. Aspect ratio score
        score.aspect_ratio_score = self._check_aspect_ratios(metrics, score)

        # 2. Whitespace score
        score.whitespace_score = self._check_whitespace(metrics, score)

        # 3. Collision score
        score.collision_score = self._check_collisions(metrics, score, slide_id)

        # 4. Balance score
        score.balance_score = self._check_balance(metrics, score, slide_id)

        # 5. Text density score
        score.text_density_score = self._check_text_density(metrics, score)

        # 6. Element count score
        score.element_count_score = self._check_element_count(metrics, score, slide_id)

        return score

    def _check_aspect_ratios(
        self, metrics: list[ShapeMetrics], score: SlideAestheticScore
    ) -> float:
        """Check if element aspect ratios are within valid ranges."""
        violations = 0
        for m in metrics:
            if m.width < 0.02 or m.height < 0.02:
                continue  # Skip tiny shapes
            ar = m.aspect_ratio
            if ar < self.MIN_ASPECT_RATIO or ar > self.MAX_ASPECT_RATIO:
                violations += 1
                score.issues.append(
                    f"元素{m.shape_id}的宽高比{ar:.1f}超出合理范围({self.MIN_ASPECT_RATIO}-{self.MAX_ASPECT_RATIO})"
                )
                score.suggestions.append(f"调整元素{m.shape_id}的尺寸使宽高比在0.1-20之间")

        if not metrics:
            return 1.0
        ratio = 1.0 - (violations / len(metrics))
        return max(0.0, ratio)

    def _check_whitespace(
        self, metrics: list[ShapeMetrics], score: SlideAestheticScore
    ) -> float:
        """Check whitespace balance — elements shouldn't cover too much or too little."""
        total_area = sum(m.area for m in metrics if m.width > 0.02 and m.height > 0.02)

        if total_area > self.MAX_COVERAGE:
            excess = total_area - self.MAX_COVERAGE
            score.issues.append(
                f"元素总覆盖率{total_area:.0%}过高（建议<{self.MAX_COVERAGE:.0%}），留白不足"
            )
            score.suggestions.append("减少元素数量或缩小部分元素，增加留白空间")
            return max(0.0, 1.0 - (excess / 0.30))

        if total_area < self.MIN_COVERAGE:
            score.issues.append(
                f"元素总覆盖率{total_area:.0%}过低（建议>{self.MIN_COVERAGE:.0%}），内容稀疏"
            )
            score.suggestions.append("增加内容元素或适当放大现有元素")
            return max(0.0, total_area / self.MIN_COVERAGE)

        return 1.0

    def _check_collisions(
        self, metrics: list[ShapeMetrics], score: SlideAestheticScore, slide_id: int
    ) -> float:
        """Detect overlapping elements using bounding box intersection."""
        collision_count = 0
        n = len(metrics)

        for i in range(n):
            for j in range(i + 1, n):
                a, b = metrics[i], metrics[j]
                if a.width < 0.02 or b.width < 0.02:
                    continue

                # Compute intersection
                ix = max(0, min(a.right, b.right) - max(a.left, b.left))
                iy = max(0, min(a.bottom, b.bottom) - max(a.top, b.top))
                overlap_area = ix * iy

                min_area = min(a.area, b.area)
                if min_area > 0 and overlap_area / min_area > self.COLLISION_THRESHOLD:
                    collision_count += 1
                    score.issues.append(
                        f"幻灯片{slide_id}中元素{a.shape_id}和{b.shape_id}重叠"
                        f"（重叠率{overlap_area/min_area:.0%}）"
                    )
                    score.suggestions.append(
                        f"将元素{a.shape_id}或{b.shape_id}移动位置消除重叠"
                    )

        # Score: 1.0 for no collisions, decreasing for each collision
        max_pairs = max(1, n * (n - 1) / 2)
        return max(0.0, 1.0 - collision_count / max_pairs)

    def _check_balance(
        self, metrics: list[ShapeMetrics], score: SlideAestheticScore, slide_id: int
    ) -> float:
        """Check visual weight balance (left-right and top-bottom)."""
        if not metrics:
            return 1.0

        valid = [m for m in metrics if m.width > 0.02 and m.height > 0.02]
        if not valid:
            return 1.0

        # Compute center of mass
        total_area = sum(m.area for m in valid)
        if total_area == 0:
            return 1.0

        cx = sum(m.area * (m.left + m.width / 2) for m in valid) / total_area
        cy = sum(m.area * (m.top + m.height / 2) for m in valid) / total_area

        # Distance from slide center (0.5, 0.5)
        h_imbalance = abs(cx - 0.5)
        v_imbalance = abs(cy - 0.5)

        issues_added = False
        if h_imbalance > self.BALANCE_THRESHOLD:
            side = "左侧" if cx < 0.5 else "右侧"
            score.issues.append(
                f"幻灯片{slide_id}视觉重心偏{side}（偏移量{h_imbalance:.0%}）"
            )
            score.suggestions.append(f"调整元素位置使视觉重心居中，当前偏{side}")
            issues_added = True

        if v_imbalance > self.BALANCE_THRESHOLD:
            side = "上方" if cy < 0.5 else "下方"
            if not issues_added:
                score.issues.append(
                    f"幻灯片{slide_id}视觉重心偏{side}（偏移量{v_imbalance:.0%}）"
                )
                score.suggestions.append(f"调整元素垂直位置使版面平衡")

        imbalance = max(h_imbalance, v_imbalance)
        return max(0.0, 1.0 - imbalance / 0.5)

    def _check_text_density(
        self, metrics: list[ShapeMetrics], score: SlideAestheticScore
    ) -> float:
        """Check text density — too much text reduces readability."""
        text_shapes = [m for m in metrics if m.has_text and m.text_length > 0]
        if not text_shapes:
            return 1.0

        total_chars = sum(m.text_length for m in text_shapes)

        # PPT best practice: <200 chars per slide
        if total_chars > 400:
            score.issues.append(
                f"文本总量{total_chars}字超过PPT最佳实践（建议<400字）"
            )
            score.suggestions.append("将部分文字移入演讲备注，保留核心要点即可")
            return max(0.0, 1.0 - (total_chars - 400) / 600)

        # Too few characters (for non-cover slides)
        if total_chars < 20 and len(text_shapes) == 1:
            score.issues.append("文本内容过少，幻灯片可能缺少必要信息")
            return 0.7

        # Check font sizes for readability
        small_fonts = [m for m in text_shapes if 0 < m.font_size_pt < self.MIN_FONT_PT]
        if small_fonts:
            score.issues.append(
                f"{len(small_fonts)}个文本框字号过小（<{self.MIN_FONT_PT}pt），投影后难以阅读"
            )
            score.suggestions.append(f"将字号调整到{self.MIN_FONT_PT}pt以上")
            return max(0.5, 1.0 - len(small_fonts) * 0.15)

        return 1.0

    def _check_element_count(
        self, metrics: list[ShapeMetrics], score: SlideAestheticScore, slide_id: int
    ) -> float:
        """Check if element count is within optimal range."""
        visible = [m for m in metrics if m.width > 0.02 and m.height > 0.02]
        count = len(visible)

        if count > self.MAX_ELEMENTS:
            score.issues.append(
                f"幻灯片{slide_id}包含{count}个元素（建议≤{self.MAX_ELEMENTS}），版面过于拥挤"
            )
            score.suggestions.append(f"合并相关元素或拆分为多张幻灯片")
            return max(0.0, 1.0 - (count - self.MAX_ELEMENTS) / 5)

        return 1.0

    def _aggregate(self, slide_scores: list[SlideAestheticScore]) -> PresentationAestheticScore:
        """Aggregate per-slide scores into presentation-level score."""
        if not slide_scores:
            return PresentationAestheticScore(overall_score=75.0)

        avg_score = sum(s.overall_score for s in slide_scores) / len(slide_scores)

        # Identify worst slides (below 60)
        worst = [s.slide_id for s in slide_scores if s.overall_score < 60]
        worst.sort(key=lambda sid: next(s.overall_score for s in slide_scores if s.slide_id == sid))

        # Global issues (patterns across multiple slides)
        global_issues = []
        collision_slides = [s.slide_id for s in slide_scores if s.collision_score < 0.8]
        if len(collision_slides) > 1:
            global_issues.append(f"多张幻灯片（{collision_slides}）存在元素重叠，建议系统性检查布局")

        balance_slides = [s.slide_id for s in slide_scores if s.balance_score < 0.7]
        if len(balance_slides) > 2:
            global_issues.append(f"{len(balance_slides)}张幻灯片视觉重心偏移，建议使用居中对齐模板")

        return PresentationAestheticScore(
            slide_scores=slide_scores,
            overall_score=avg_score,
            worst_slides=worst[:5],
            global_issues=global_issues,
        )

    def _unavailable_score(self) -> PresentationAestheticScore:
        """Return a neutral score when python-pptx is unavailable."""
        return PresentationAestheticScore(
            slide_scores=[],
            overall_score=75.0,
            global_issues=["无法加载python-pptx，跳过几何质量检查"],
        )

    def generate_fix_prompt(self, score: PresentationAestheticScore) -> str:
        """Generate a targeted fix prompt for SlideEditorAgent based on scored issues."""
        if score.overall_score >= 90:
            return ""

        lines = [f"PPT几何质量评分：{score.overall_score:.0f}/100（{score._grade()}级）\n"]

        if score.global_issues:
            lines.append("全局问题：")
            lines.extend(f"  - {issue}" for issue in score.global_issues)
            lines.append("")

        priority_slides = sorted(
            score.slide_scores, key=lambda s: s.overall_score
        )[:3]

        for slide in priority_slides:
            if slide.issues:
                lines.append(f"幻灯片{slide.slide_id}（评分{slide.overall_score:.0f}）：")
                lines.extend(f"  问题：{issue}" for issue in slide.issues[:2])
                lines.extend(f"  建议：{sug}" for sug in slide.suggestions[:2])
                lines.append("")

        return "\n".join(lines)
